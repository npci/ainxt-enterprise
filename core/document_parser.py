# SPDX-License-Identifier: Apache-2.0
# ============================================================
# DOCUMENT PARSER — unified multimodal file extraction
#
# Routes by file_type → returns clean extracted text.
# Outputs are truncated only if they exceed _MAX_CHARS (2_000_000 chars).
#
# Supported types:
#   pdf, docx, xlsx, xls, csv, html, rtf, txt, json,
#   png, jpg, jpeg, gif, webp, bmp (→ Gemini Vision)
# ============================================================

import base64
import json
import os
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

from core.logger import logger

_MAX_CHARS = 2_000_000   # raised from 100_000 — large PDFs (200+ pages) must be fully indexed


# ============================================================
# INDIVIDUAL PARSERS
# ============================================================

def parse_pdf(path: str) -> str:
    """
    Convert PDF to Markdown.

    Markdown extraction ladder, best fidelity first:
      1. markitdown (MIT) — via core.pdf_backend.to_markdown(). No heading
         inference, no batching, no header/footer strip, so structure is
         coarser than the AGPL-3.0 pymupdf4llm engine this replaced, but text
         content is not lost. See core/pdf_backend/pdfium_backend.py.
      2. _parse_pdf_plain — always available via core.pdf_backend; also runs
         the hybrid OCR fallback for scanned/mixed PDFs.
    """
    try:
        from core import pdf_backend as _be
        md = _be.to_markdown(path)
        if md and md.strip():
            logger.debug("parse_pdf: markitdown produced %d chars", len(md))
            return md.strip()
    except Exception as e:
        logger.debug(f"parse_pdf: markitdown unavailable/failed ({e}) — falling back to plain-text")
    return _parse_pdf_plain(path)


def _parse_pdf_plain(path: str) -> str:
    """Plain-text fallback using core.pdf_backend page.get_text(). Filters blank/image-only pages.

    If every page is empty (scanned / image-only), delegates to the hybrid
    OCR + table extractor in ``core.pdf_ocr``. This keeps born-digital PDFs
    on the zero-OCR fast path while making scanned PDFs work end-to-end.
    """
    try:
        from core import pdf_backend as fitz
        with fitz.open(path) as doc:
            pages = [page.get_text().strip() for page in doc]
        # Drop blank pages (image-only, empty, or whitespace-only pages)
        nonblank = [p for p in pages if p]
        if nonblank:
            return "\n\n".join(nonblank)
        # All pages blank → try hybrid OCR extractor before giving up.
        return _hybrid_ocr_fallback(path)
    except Exception as e:
        logger.error(f"parse_pdf plain fallback failed: {e}")
        return f"[PDF parse error: {e}]"


def _hybrid_ocr_fallback(path: str) -> str:
    """Run the hybrid extractor (born-digital text + tables + scanned-page OCR).

    Returns "" if OCR is unavailable so callers can decide whether to surface
    a user-facing error. Any extractor exception is swallowed and logged —
    we never want OCR to be the reason a parse call raises.
    """
    try:
        from core.pdf_ocr import extract_pdf
        result = extract_pdf(path)
        text = (result.get("text") or "").strip()
        if text:
            logger.info(
                "pdf_ocr: hybrid extractor recovered %d chars (source=%s, engine=%s)",
                len(text), result.get("source"), result.get("ocr_engine"),
            )
            return text
        warnings = result.get("warnings") or []
        if warnings:
            logger.warning("pdf_ocr: extractor returned empty. warnings=%s", warnings[:5])
        return ""
    except Exception as exc:  # pragma: no cover — defensive
        logger.warning("pdf_ocr: hybrid fallback raised: %s", exc)
        return ""


def parse_docx(path: str) -> str:
    """
    Convert DOCX to Markdown preserving heading levels, tables, and lists.
    Heading 1 → #, Heading 2 → ##, Heading 3 → ###, etc.
    Tables → GitHub-flavoured Markdown pipe tables.
    Lists → - bullet items.
    """
    try:
        import docx
        from docx.oxml.ns import qn

        doc   = docx.Document(path)
        lines = []

        # Map Word style names → Markdown heading prefix
        _HEADING = {
            "heading 1": "#", "heading 2": "##", "heading 3": "###",
            "heading 4": "####", "heading 5": "#####", "heading 6": "######",
            "title": "#", "subtitle": "##",
        }

        def _table_to_md(table) -> str:
            rows = []
            for i, row in enumerate(table.rows):
                cells = [c.text.replace("\n", " ").strip() for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
                if i == 0:  # header separator
                    rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
            return "\n".join(rows)

        def _is_bold_heading(para) -> bool:
            """
            Detect paragraphs that act as section headings via bold formatting
            rather than Word heading styles. This handles documents (like policy
            manuals) where authors use bold Normal-style text instead of Heading 1/2/3.

            A paragraph is treated as a bold heading when ALL of the following hold:
              1. It uses a non-heading style (Normal, Body Text, etc.)
              2. Every non-whitespace run is bold
              3. The text is short (≤ 120 chars) — headings are never long sentences
              4. The text does not end with sentence-ending punctuation (. ! ?)
                 — body sentences end with punctuation; headings typically don't
            """
            style = para.style.name.lower() if para.style else ""
            if style in _HEADING:
                return False  # already handled as a proper heading style
            text = para.text.strip()
            if not text or len(text) > 120:
                return False
            if text[-1] in ".!?,:;":
                return False
            # Check that every non-empty run is bold
            non_empty_runs = [r for r in para.runs if r.text.strip()]
            if not non_empty_runs:
                return False
            return all(r.bold for r in non_empty_runs)

        def _bold_heading_level(para) -> str:
            """
            Map a bold-heading paragraph to a Markdown heading level.
            Uses font size as a proxy for hierarchy:
              ≥ 14pt  → ## (level 2 — major section)
              < 14pt or no size → ### (level 3 — sub-section)
            """
            sizes = [r.font.size.pt for r in para.runs if r.font.size]
            max_size = max(sizes) if sizes else 0
            return "##" if max_size >= 14 else "###"

        def _emit_paragraph(para) -> None:
            text = para.text.strip()
            if not text:
                lines.append("")
                return
            style = para.style.name.lower() if para.style else ""
            if style in _HEADING:
                lines.append(f"{_HEADING[style]} {text}")
            elif style in ("list paragraph", "list bullet", "list number"):
                lines.append(f"- {text}")
            elif _is_bold_heading(para):
                # Bold Normal-style paragraph acting as a section heading.
                # Emit as a Markdown heading so _make_section_map() can
                # detect it and build a proper section_map for Coverage tier.
                hprefix = _bold_heading_level(para)
                lines.append(f"{hprefix} {text}")
            else:
                lines.append(text)

        # Word's "Content Control" feature (Structured Document Tags, <w:sdt>)
        # wraps its content in <w:sdt><w:sdtContent>...actual <w:p>/<w:tbl>
        # elements...</w:sdtContent></w:sdt>. python-docx's own doc.paragraphs
        # and doc.tables ONLY see block-level <w:p>/<w:tbl> elements — anything
        # nested inside an <w:sdt> is invisible to them. Templated documents
        # (policy/finance manuals, government forms, DMS-managed docs) commonly
        # use content controls for version blocks, approval fields, and
        # placeholder text, so a doc built entirely (or mostly) from content
        # controls was silently parsed as empty — no exception, no sentinel,
        # just "" — which surfaced upstream as "File parsed but contained no
        # readable text." with no indication of the real cause.
        #
        # Fix: walk the body recursively so <w:sdt>/<w:sdtContent> (and any
        # other unrecognised wrapper tag) is descended into rather than
        # skipped, while <w:p>/<w:tbl> keep being handled exactly as before —
        # this changes nothing for documents that don't use content controls.
        _SDT_WRAPPER_TAGS = {"sdt", "sdtContent"}
        _MAX_SDT_DEPTH = 20  # guard against pathological/malformed nesting

        def _walk(elements, depth=0) -> None:
            if depth > _MAX_SDT_DEPTH:
                return
            for block in elements:
                tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

                if tag == "tbl":
                    # Find the corresponding Table object. doc.tables only
                    # lists TOP-LEVEL tables, so a table nested inside an sdt
                    # (or inside another table's cell) won't be found there —
                    # construct the Table wrapper directly from the element
                    # instead, exactly as python-docx does internally.
                    tbl = next((t for t in doc.tables if t._tbl is block), None)
                    if tbl is None:
                        from docx.table import Table as _Table
                        tbl = _Table(block, doc)
                    lines.append(_table_to_md(tbl))
                    lines.append("")

                elif tag == "p":
                    # Same reasoning as above — construct the Paragraph
                    # wrapper directly rather than requiring it be present in
                    # doc.paragraphs (which excludes sdt-nested paragraphs).
                    para = next((p for p in doc.paragraphs if p._p is block), None)
                    if para is None:
                        from docx.text.paragraph import Paragraph as _Paragraph
                        para = _Paragraph(block, doc)
                    _emit_paragraph(para)

                elif tag in _SDT_WRAPPER_TAGS:
                    # Descend into the content-control wrapper. <w:sdt> holds
                    # exactly one <w:sdtPr> (properties — skipped, not content)
                    # and one <w:sdtContent> (the actual body elements); walking
                    # all children and letting the tag check filter is simpler
                    # and robust to either being passed directly.
                    _walk(list(block), depth + 1)
                # Any other tag (w:sdtPr, w:bookmarkStart, w:proofErr, sectPr,
                # etc.) carries no document text and is silently skipped, same
                # as before this fix.

        _walk(doc.element.body)

        return "\n".join(lines)

    except ImportError:
        return "[DOCX parsing unavailable — install python-docx]"
    except Exception as e:
        logger.error(f"parse_docx failed: {e}")
        return f"[DOCX parse error: {e}]"


def parse_excel(path: str) -> str:
    """
    Convert Excel to Markdown tables using tabulate.
    Multi-sheet workbooks: each sheet becomes a ## heading + table.
    """
    try:
        import pandas as pd
        try:
            from tabulate import tabulate
            _tabulate = lambda df: tabulate(df, headers="keys", tablefmt="github", showindex=False)
        except ImportError:
            _tabulate = lambda df: df.to_string(index=False)

        xl = pd.ExcelFile(path)
        parts = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            if df.empty:
                continue
            if len(xl.sheet_names) > 1:
                parts.append(f"## {sheet}\n")
            parts.append(_tabulate(df))
        return "\n\n".join(parts) if parts else "[Excel file is empty]"
    except ImportError:
        return "[Excel parsing unavailable — install pandas + openpyxl]"
    except Exception as e:
        logger.error(f"parse_excel failed: {e}")
        return f"[Excel parse error: {e}]"


def parse_csv(path: str) -> str:
    """Convert CSV to a GitHub-flavoured Markdown table using tabulate."""
    try:
        import pandas as pd
        try:
            from tabulate import tabulate
            df = pd.read_csv(path)
            return tabulate(df, headers="keys", tablefmt="github", showindex=False)
        except ImportError:
            df = pd.read_csv(path)
            return df.to_string(index=False)
    except ImportError:
        return "[CSV parsing unavailable — install pandas]"
    except Exception as e:
        logger.error(f"parse_csv failed: {e}")
        return f"[CSV parse error: {e}]"


def parse_html(path: str) -> str:
    try:
        from bs4 import BeautifulSoup
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "html.parser")
        # Remove noise tags before text extraction — get_text() would otherwise
        # include raw JS, CSS, and <head> metadata in the extracted content,
        # polluting every RAG chunk with irrelevant code.
        for tag in soup(["script", "style", "head", "meta", "noscript"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        return "[HTML parsing unavailable — install beautifulsoup4]"
    except Exception as e:
        logger.error(f"parse_html failed: {e}")
        return f"[HTML parse error: {e}]"


def parse_rtf(path: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return rtf_to_text(f.read())
    except ImportError:
        return "[RTF parsing unavailable — install striprtf]"
    except Exception as e:
        logger.error(f"parse_rtf failed: {e}")
        return f"[RTF parse error: {e}]"


def parse_image(path: str, filename: str) -> str:
    """Send image to vision model via LLM proxy; fallback to placeholder if unavailable.

    Routes through LLM_PROXY_URL → POST /llm/generate-image (Gemini vision).
    This is the same proxy path used by /ask/image and the image-intent routing
    block in gateway.py — all Gemini calls go through the proxy on web02.
    """
    _VISION_PROMPT = (
        "Describe this image in detail, including all visible text, numbers, "
        "charts, diagrams, and visual elements."
    )
    try:
        import base64 as _b64
        import urllib.request as _ur
        import json as _json
        _proxy = os.environ.get("LLM_PROXY_URL", "").rstrip("/")
        if not _proxy:
            raise ValueError("LLM_PROXY_URL not set")
        with open(path, "rb") as f:
            image_data = f.read()
        ext = Path(path).suffix.lstrip(".").lower()
        mime_map = {
            "jpg": "image/jpeg", "jpeg": "image/jpeg",
            "png": "image/png", "gif": "image/gif",
            "webp": "image/webp", "bmp": "image/bmp",
        }
        mime = mime_map.get(ext, "image/jpeg")
        _img_b64 = _b64.b64encode(image_data).decode()
        _payload = _json.dumps({
            "provider":  "gemini",
            "prompt":    _VISION_PROMPT,
            "image_b64": _img_b64,
            "mime_type": mime,
        }).encode()
        _req = _ur.Request(
            f"{_proxy}/llm/generate-image",
            data=_payload,
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        import contextlib
        with contextlib.closing(_ur.urlopen(_req, timeout=60)) as _resp:
            _data = _json.loads(_resp.read())
        _text = (_data.get("text") or "").strip()
        if _text:
            logger.info(f"parse_image: proxy vision OK for {filename!r} ({len(_text)} chars)")
            return _text
        logger.warning(f"parse_image: proxy returned empty text for {filename!r}")
        return f"[Image: {filename} — vision unavailable]"
    except Exception as e:
        logger.warning(f"parse_image: vision unavailable ({e}); using fallback")
        return f"[Image: {filename} — vision unavailable]"


def parse_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"[TXT read error: {e}]"


def parse_json(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            data = json.load(f)
        return json.dumps(data, indent=2)
    except Exception as e:
        return f"[JSON parse error: {e}]"


def parse_pptx(path: str) -> str:
    """Extract text from PowerPoint (.pptx) slide by slide using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            # Slide title via layout/placeholder
            title_ph = slide.shapes.title
            if title_ph and title_ph.has_text_frame:
                title_text = title_ph.text_frame.text.strip()
                if title_text:
                    parts.append(f"## {title_text}")
            for shape in slide.shapes:
                if shape == title_ph:
                    continue
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(f"- {line}" if para.level > 0 else line)
                elif shape.has_table:
                    tbl = shape.table
                    rows = []
                    for r_idx, row in enumerate(tbl.rows):
                        cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                        if r_idx == 0:
                            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    parts.append("\n".join(rows))
            if parts:
                slides.append(f"### Slide {i}\n" + "\n".join(parts))
        return "\n\n".join(slides) if slides else "[Presentation has no text content]"
    except ImportError:
        return "[PPTX parsing unavailable — install python-pptx]"
    except Exception as e:
        logger.error(f"parse_pptx failed: {e}")
        return f"[PPTX parse error: {e}]"


# ============================================================
# DISPATCHER
# ============================================================

_IMAGE_TYPES = {"png", "jpg", "jpeg", "gif", "webp", "bmp"}

# Formats Docling can handle. Kept in sync with core.docling_parser._SUPPORTED
# so the dispatcher never hands a non-supported format to Docling.
_DOCLING_FORMATS = {"pdf", "docx", "html", "htm", "pptx"}


def _try_docling(path: str, file_type: str) -> Optional[str]:
    """
    Run Docling for `file_type`.

    Two execution paths — chosen by whether PARSE_SVC_URL is configured:

    Remote path (PARSE_SVC_URL is set):
        Sends the file bytes to the embed server's POST /parse endpoint via
        httpx.  Docling + PaddleOCR run on the embed server, not in this
        gateway process.  This offloads heavy ML models from the gateway
        worker — exactly mirroring how embedding was offloaded to EMBED_SVC_URL.

    Local path (PARSE_SVC_URL is empty):
        Runs Docling in-process (legacy behavior, unchanged).
        USE_DOCLING_PARSER flag still controls activation.

    Returns None to signal "fall back to the legacy parser" — the caller
    MUST treat None as fallback, never as empty content.

    Failure semantics: raises an exception on any failure — timeout, network
    error, empty content, or model error. The caller (activate_doc) must treat
    this as a hard failure and NOT fall back to the legacy parser, because
    legacy-parsed embeddings for Docling-supported formats produce incorrect
    chunking quality and mislead RAG retrieval.

    Returns the parsed markdown string on success (always non-empty).
    Raises RuntimeError with a descriptive message on any failure.
    """
    ft = (file_type or "").lower().strip(".")

    # ── Remote path: delegate to parse service ────────────────────────────────
    try:
        from core.config import PARSE_SVC_URL as _PARSE_SVC_URL, PARSE_SVC_TIMEOUT as _PARSE_TIMEOUT
    except Exception:
        _PARSE_SVC_URL = ""
        _PARSE_TIMEOUT = 1800.0   # matches core/config.py default — 30 min for large files

    if _PARSE_SVC_URL:
        import base64 as _b64
        import httpx as _httpx

        if ft not in _DOCLING_FORMATS:
            # Format not supported by Docling — not a failure, just not applicable
            return None

        with open(path, "rb") as _fh:
            _file_bytes = _fh.read()

        _fname = os.path.basename(path)
        logger.info(
            f"_try_docling: calling parse service "
            f"url={_PARSE_SVC_URL} file='{_fname}' "
            f"ext={ft} size={len(_file_bytes):,}B timeout={_PARSE_TIMEOUT}s"
        )

        _timeout = _httpx.Timeout(
            connect=10.0,
            read=_PARSE_TIMEOUT,
            write=30.0,
            pool=10.0,
        )
        logger.info(
            f"_try_docling: before http call for /parse "
            f"file='{_fname}' ext={ft} size={len(_file_bytes):,}B"
        )

        try:
            _resp = _httpx.post(
                f"{_PARSE_SVC_URL.rstrip('/')}/parse",
                json={
                    "file_bytes_b64": _b64.b64encode(_file_bytes).decode(),
                    "filename":       _fname,
                    "file_type":      ft,
                },
                timeout=_timeout,
            )
        except _httpx.ConnectError as e:
            logger.error(
                f"_try_docling: parse service unreachable "
                f"url={_PARSE_SVC_URL} file='{_fname}' error='{e}'"
            )
            # User-facing message: no service URL, no model name, no exception
            # class. The logger.error() above retains all of that for ops.
            raise RuntimeError(
                "The document processing service is currently unavailable. "
                "Please try again in a few minutes, or contact support if the "
                "problem continues."
            )
        except _httpx.ReadTimeout:
            logger.error(
                f"_try_docling: parse service read timeout after {_PARSE_TIMEOUT}s "
                f"file='{_fname}' ext={ft} — "
                f"increase PARSE_SVC_TIMEOUT in .env for large files"
            )
            # User-facing: keep the filename and the elapsed minutes (both
            # actionable), drop the env-var tuning hint (ops-only) and the
            # engine name. Full detail stays in the logger.error() above.
            raise RuntimeError(
                f"Processing '{_fname}' took longer than "
                f"{max(1, int(_PARSE_TIMEOUT // 60))} minute(s) and was stopped. "
                f"The document may be very large or complex. Please try again, "
                f"or split it into smaller files."
            )
        except Exception as e:
            logger.error(
                f"_try_docling: parse service HTTP call failed "
                f"file='{_fname}' error='{type(e).__name__}: {e}'"
            )
            raise RuntimeError(
                f"Could not process '{_fname}' because the document processing "
                f"service could not be reached. Please try again, or contact "
                f"support if the problem continues."
            )

        # HTTP 422 = page-level conversion failure. The parse service puts the
        # user-facing message (exact failed page ranges + total page count) in
        # the `detail` field as a plain STRING. Surface that verbatim so
        # activate_doc() can store it in knowledge_docs.parse_error and the user
        # sees precisely which pages could not be extracted.
        #
        # Caveat: FastAPI ALSO returns 422 for its own request-validation
        # failures, where `detail` is a list of dicts, e.g.
        #     [{"loc": ["body", "file_bytes_b64"], "msg": "field required", ...}]
        # Stringifying that would expose internal field names and schema detail
        # in the Request Status tab. So only a str detail is trusted as
        # user-facing; anything else (list/dict/empty) becomes a generic message.
        # The raw body is logged either way for diagnosis.
        if _resp.status_code == 422:
            try:
                _detail = _resp.json().get("detail")
            except Exception:
                _detail = None

            logger.error(
                f"_try_docling: parse service returned HTTP 422 "
                f"file='{_fname}' detail={_detail!r} body='{_resp.text[:500]}'"
            )

            if isinstance(_detail, str) and _detail.strip():
                # Genuine page-conversion failure message from the parse service.
                raise RuntimeError(_detail.strip())

            # Request-validation error or unexpected shape — never surface the
            # raw payload (field names, schema internals, JSON structure).
            raise RuntimeError(
                f"Could not process '{_fname}'. The document could not be read "
                f"by the document processing service. Please verify the file "
                f"opens correctly and re-upload it."
            )

        try:
            _resp.raise_for_status()
        except Exception as e:
            logger.error(
                f"_try_docling: parse service returned HTTP {_resp.status_code} "
                f"file='{_fname}' body='{_resp.text[:200]}'"
            )
            # Do NOT surface the raw response body — it can contain tracebacks,
            # module paths and model names. Logged in full just above.
            raise RuntimeError(
                f"Could not process '{_fname}'. The document processing service "
                f"reported an unexpected error. Please try again, or contact "
                f"support if the problem continues."
            )

        _content = _resp.json().get("content", "")
        logger.info(
            f"_try_docling: after http call for /parse "
            f"file='{_fname}' http_status={_resp.status_code} chars={len(_content):,}"
        )

        if not _content or not _content.strip():
            logger.error(
                f"_try_docling: parse service returned empty content "
                f"file='{_fname}' ext={ft}"
            )
            raise RuntimeError(
                f"No readable text could be extracted from '{_fname}'. "
                f"The file may be corrupt, password-protected, or in an "
                f"unsupported format. Please check the file and re-upload."
            )

        logger.info(
            f"_try_docling: parse service OK "
            f"file='{_fname}' chars={len(_content):,}"
        )
        return _content

    # ── Local path: in-process Docling ────────────────────────────────────────
    try:
        from core import docling_parser as _dp
        if not _dp.is_active():
            return None
        if not _dp.supports(file_type):
            return None
        _result = _dp.parse(path, file_type)
        if not _result or not _result.strip():
            logger.error(
                f"_try_docling: in-process Docling returned empty content "
                f"file='{os.path.basename(path)}' ext={ft}"
            )
            raise RuntimeError(
                f"No readable text could be extracted from "
                f"'{os.path.basename(path)}'. The file may be corrupt, "
                f"password-protected, or in an unsupported format. "
                f"Please check the file and re-upload."
            )
        return _result
    except RuntimeError:
        # Re-raises both plain RuntimeError and PageConversionError (which
        # inherits from RuntimeError).  PageConversionError carries the exact
        # failed page ranges and must never be swallowed here — it propagates
        # to activate_doc() which rolls back the document to PENDING_APPROVAL.
        # Legacy parsing must NOT silently replace a partial Docling result.
        raise
    except Exception as e:
        logger.error(
            f"_try_docling: in-process Docling failed "
            f"file='{os.path.basename(path)}' ext={ft} error='{type(e).__name__}: {e}'"
        )
        # Generic catch-all. The exception type and message are logged above;
        # the user only needs to know which file failed and what to do next.
        raise RuntimeError(
            f"Could not process '{os.path.basename(path)}'. "
            f"Please verify the file opens correctly and re-upload it, "
            f"or contact support if the problem continues."
        )


def parse_file(path: str, file_type: str, filename: str, skip_docling: bool = False) -> str:
    """
    Route to the correct parser by file_type.
    Returns extracted text, truncated to _MAX_CHARS.

    When USE_DOCLING_PARSER=1 (or PARSE_SVC_URL is set) and file_type is in
    _DOCLING_FORMATS, Docling is tried first — UNLESS skip_docling=True.

    skip_docling=True: bypass Docling entirely and go straight to the legacy
    parser (markitdown, python-docx, etc.).  Used at upload time so Docling
    only runs after a document is approved (inside activate_doc()).

    On any Docling failure (import error, model init, conversion error, empty
    output) we fall back to the per-format legacy parser so the upload never
    breaks because of the flag.
    """
    ft = file_type.lower().strip(".")

    # Docling fast-path — skipped when skip_docling=True (deferred to activate_doc).
    # Returns None to signal fallback. Keeps the legacy branches below
    # untouched so flag-off behavior is bit-for-bit identical.
    if not skip_docling and ft in _DOCLING_FORMATS:
        _docling_md = _try_docling(path, ft)
        if _docling_md is not None:
            text = _docling_md
            if len(text) > _MAX_CHARS:
                logger.warning(
                    f"parse_file(docling): text exceeds {_MAX_CHARS} chars "
                    f"({len(text)} chars) — truncating."
                )
                text = text[:_MAX_CHARS]
            return text

    if ft == "pdf":
        text = parse_pdf(path)
    elif ft == "docx":
        text = parse_docx(path)
    elif ft == "pptx":
        text = parse_pptx(path)
    elif ft == "ppt":
        # python-pptx only reads OOXML (.pptx / ZIP-based) files.
        # Legacy OLE2 binary .ppt (PowerPoint 97-2003) cannot be parsed.
        # Return a clear message so the stored content is not a silent error string.
        text = "[Legacy .ppt format is not supported for text extraction. Please convert to .pptx and re-upload.]"
    elif ft in ("xlsx", "xls"):
        text = parse_excel(path)
    elif ft == "csv":
        text = parse_csv(path)
    elif ft in ("html", "htm"):
        text = parse_html(path)
    elif ft == "rtf":
        text = parse_rtf(path)
    elif ft in _IMAGE_TYPES:
        text = parse_image(path, filename)
    elif ft == "json":
        text = parse_json(path)
    elif ft == "xml":
        # XML files are plain text — read them as-is so the LLM sees the
        # full structure (tags + content).  Same path as .txt/.md.
        text = parse_txt(path)
    elif ft in ("txt", "md"):
        # .md is plain text — Markdown syntax is preserved as-is, which is
        # correct for RAG chunking (the chunker is Markdown heading-aware).
        text = parse_txt(path)
    else:
        # Attempt plain text for unknown types
        text = parse_txt(path)

    if len(text) > _MAX_CHARS:
        logger.warning(
            f"parse_file: text exceeds {_MAX_CHARS} chars ({len(text)} chars) — "
            f"truncating. Consider raising _MAX_CHARS for this document type."
        )
        text = text[:_MAX_CHARS]

    return text


def parse_file_structured(path: str, file_type: str, filename: str, skip_docling: bool = False) -> dict:
    """
    Parse file and return a standardised enterprise document object.

    skip_docling=True: bypass Docling and use legacy parsers only.
    Used at KB upload time — Docling is deferred to activate_doc() so
    wasted parse calls never happen for docs deleted before approval.

    Returns:
        {
            "type":      "pdf|docx|excel|image|...",
            "content":   "<extracted text>",
            "metadata":  {"filename": ..., "size_bytes": ..., "pages": ...},
            "source":    "<filename>",
            "timestamp": "<ISO-8601 UTC>"
        }
    """
    content = parse_file(path, file_type, filename, skip_docling=skip_docling)
    size_bytes = 0
    try:
        size_bytes = Path(path).stat().st_size
    except Exception:
        pass

    ft = file_type.lower().strip(".")
    doc_type = "image" if ft in _IMAGE_TYPES else ft

    # Attempt page count for PDFs
    pages = None
    if ft == "pdf":
        try:
            from core import pdf_backend as fitz
            with fitz.open(path) as doc:
                pages = len(doc)
        except Exception:
            pass

    metadata: dict = {"filename": filename, "size_bytes": size_bytes}
    if pages is not None:
        metadata["pages"] = pages

    return {
        "type":      doc_type,
        "content":   content,
        "metadata":  metadata,
        "source":    filename,
        "timestamp": datetime.now(timezone.utc).isoformat(),
    }