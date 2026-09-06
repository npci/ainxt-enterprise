# SPDX-License-Identifier: MIT
"""
Document generation utilities.
Produces .docx, .pptx, .pdf, .xlsx, .txt, .md from structured sections.

sections format (legacy):
  [{"heading": str, "content": str, "bullets": list[str], "level": int}]

PPTX rich format (from _build_pptx_prompt):
  [{"slide_type": str, "heading": str, "key_message": str, "bullets": list,
    "stats": list, "quote": str, "attribution": str,
    "two_col_left": dict, "two_col_right": dict,
    "icon": str, "image_prompt": str, "_image_bytes": bytes,
    "speaker_notes": str}]
"""
import io
import os
import re

from core import config
from core.logger import logger

TEMPLATE_PATH = config.PPT_TEMPLATE_PATH

FORMAT_EXTENSIONS = {
    "docx":         "docx",
    "word":         "docx",
    "doc":          "docx",
    "pptx":         "pptx",
    "ppt":          "pptx",
    "powerpoint":   "pptx",
    "presentation": "pptx",
    "slides":       "pptx",
    "pdf":          "pdf",
    "xlsx":         "xlsx",
    "excel":        "xlsx",
    "spreadsheet":  "xlsx",
    "txt":          "txt",
    "text":         "txt",
    "md":           "md",
    "markdown":     "md",
    "csv":          "csv",
}

MIME_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf":  "application/pdf",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "txt":  "text/plain; charset=utf-8",
    "md":   "text/markdown; charset=utf-8",
    "csv":  "text/csv; charset=utf-8",
}

# ── PPTX Theme catalogue ─────────────────────────────────────────────────────
# Plain dicts so they are importable even when python-pptx is not installed.
# Gradient stops: [(pos_pct 0-100, hex_color_no_hash), ...]
PPTX_THEMES = {
    "dark_executive": {
        "id":          "dark_executive",
        "name":        "Dark Executive",
        "description": "Deep navy gradient — classic enterprise look",
        "swatch":      "#003366",
        "preview":     "dark",
        "dark_stops":  [(0, "060D1A"), (50, "0A1A2E"), (100, "003366")],
        "light_stops": [(0, "F0F4F8"), (100, "E2EAF5")],
        "dark_angle":  135,   # pptx clockwise-from-East: NW→SE diagonal
        "light_angle": 90,    # pptx clockwise-from-East: top→bottom
    },
    "light_modern": {
        "id":          "light_modern",
        "name":        "Light Modern",
        "description": "Clean white gradient with deep-blue accents — crisp minimal",
        "swatch":      "#1A73E8",
        "preview":     "light",
        "dark_stops":  [(0, "1A2744"), (100, "16213E")],
        "light_stops": [(0, "FFFFFF"), (100, "F0F4FF")],
        "dark_angle":  90,
        "light_angle": 135,
    },
    "vibrant_tech": {
        "id":          "vibrant_tech",
        "name":        "Vibrant Tech",
        "description": "Navy-to-teal diagonal gradient — dynamic high-energy feel",
        "swatch":      "#005B8E",
        "preview":     "dark",
        "dark_stops":  [(0, "050F1E"), (50, "001833"), (100, "003355")],
        "light_stops": [(0, "F0F8FF"), (100, "E0F0FF")],
        "dark_angle":  120,
        "light_angle": 90,
    },
}

# ── DOMAIN COLOUR PALETTE SYSTEM ─────────────────────────────────────────────
#
# Three-axis colour decision framework:
#   Domain/Industry → Emotional Register → Institutional Weight
#
# Each palette: primary, accent, light_bg (hex strings, no #)
# Body text is always 333333; secondary text 888888.
#
DOMAIN_PALETTES = {
    # ── Finance / Government ──────────────────────────────────────────────────
    "government":      {"primary": "1E3A5F", "accent": "2196F3", "light_bg": "E8F4FD"},
    "regulatory":      {"primary": "1E3A5F", "accent": "2196F3", "light_bg": "E8F4FD"},
    "banking":         {"primary": "0D2137", "accent": "1565C0", "light_bg": "E3EEF9"},
    "finance":         {"primary": "0D2137", "accent": "1565C0", "light_bg": "E3EEF9"},
    "payments":        {"primary": "1A237E", "accent": "3F51B5", "light_bg": "E8EAF6"},
    "fintech":         {"primary": "1A237E", "accent": "3F51B5", "light_bg": "E8EAF6"},
    "upi":             {"primary": "1A237E", "accent": "3F51B5", "light_bg": "E8EAF6"},
    # ── Healthcare ────────────────────────────────────────────────────────────
    "healthcare":      {"primary": "00695C", "accent": "009688", "light_bg": "E0F2F1"},
    "medical":         {"primary": "00695C", "accent": "009688", "light_bg": "E0F2F1"},
    # ── Technology ────────────────────────────────────────────────────────────
    "ai":              {"primary": "1B4332", "accent": "2D6A4F", "light_bg": "E9F5EE"},
    "ml":              {"primary": "1B4332", "accent": "2D6A4F", "light_bg": "E9F5EE"},
    "technology":      {"primary": "1B4332", "accent": "2D6A4F", "light_bg": "E9F5EE"},
    "cybersecurity":   {"primary": "1A1A2E", "accent": "E94560", "light_bg": "F8E8EC"},
    "security":        {"primary": "1A1A2E", "accent": "E94560", "light_bg": "F8E8EC"},
    "infrastructure":  {"primary": "263238", "accent": "546E7A", "light_bg": "ECEFF1"},
    "cloud":           {"primary": "263238", "accent": "546E7A", "light_bg": "ECEFF1"},
    # ── Legal / Compliance ────────────────────────────────────────────────────
    "legal":           {"primary": "212121", "accent": "424242", "light_bg": "F5F5F5"},
    "compliance":      {"primary": "212121", "accent": "424242", "light_bg": "F5F5F5"},
    # ── Education / Research ──────────────────────────────────────────────────
    "education":       {"primary": "1A3C5E", "accent": "2E86AB", "light_bg": "E1F0F8"},
    "research":        {"primary": "1A3C5E", "accent": "2E86AB", "light_bg": "E1F0F8"},
    # ── Cultural / Heritage ───────────────────────────────────────────────────
    "cultural":        {"primary": "7B0D1E", "accent": "C0392B", "light_bg": "FDF0F0"},
    "heritage":        {"primary": "7B0D1E", "accent": "C0392B", "light_bg": "FDF0F0"},
    # ── Sports / Analytics ────────────────────────────────────────────────────
    "sports":          {"primary": "003087", "accent": "0057E0", "light_bg": "EEF2FF"},
    "analytics":       {"primary": "003087", "accent": "0057E0", "light_bg": "EEF2FF"},
    "ipl":             {"primary": "003087", "accent": "0057E0", "light_bg": "EEF2FF"},
    # ── Sustainability / ESG ──────────────────────────────────────────────────
    "sustainability":  {"primary": "2E7D32", "accent": "43A047", "light_bg": "E8F5E9"},
    "esg":             {"primary": "2E7D32", "accent": "43A047", "light_bg": "E8F5E9"},
    "environment":     {"primary": "2E7D32", "accent": "43A047", "light_bg": "E8F5E9"},
    # ── Retail / E-commerce ───────────────────────────────────────────────────
    "retail":          {"primary": "E65100", "accent": "F57C00", "light_bg": "FFF3E0"},
    "ecommerce":       {"primary": "E65100", "accent": "F57C00", "light_bg": "FFF3E0"},
    # ── Luxury / Premium ──────────────────────────────────────────────────────
    "luxury":          {"primary": "2C1810", "accent": "8D6E63", "light_bg": "EFEBE9"},
    "premium":         {"primary": "2C1810", "accent": "8D6E63", "light_bg": "EFEBE9"},
    # ── Startups / Innovation ─────────────────────────────────────────────────
    "startup":         {"primary": "4A148C", "accent": "7B1FA2", "light_bg": "F3E5F5"},
    "innovation":      {"primary": "4A148C", "accent": "7B1FA2", "light_bg": "F3E5F5"},
    # ── Media / Entertainment ─────────────────────────────────────────────────
    "media":           {"primary": "B71C1C", "accent": "E53935", "light_bg": "FFEBEE"},
    "entertainment":   {"primary": "B71C1C", "accent": "E53935", "light_bg": "FFEBEE"},
    # ── Travel / Tourism ──────────────────────────────────────────────────────
    "travel":          {"primary": "006064", "accent": "00838F", "light_bg": "E0F7FA"},
    "tourism":         {"primary": "006064", "accent": "00838F", "light_bg": "E0F7FA"},
    # ── Food / Agriculture ────────────────────────────────────────────────────
    "food":            {"primary": "33691E", "accent": "558B2F", "light_bg": "F1F8E9"},
    "agriculture":     {"primary": "33691E", "accent": "558B2F", "light_bg": "F1F8E9"},
    # ── HR / People ───────────────────────────────────────────────────────────
    "hr":              {"primary": "AD1457", "accent": "D81B60", "light_bg": "FCE4EC"},
    "people":          {"primary": "AD1457", "accent": "D81B60", "light_bg": "FCE4EC"},
    # ── Executive / C-Suite ───────────────────────────────────────────────────
    "executive":       {"primary": "1C1C1C", "accent": "B8960C", "light_bg": "FAFAFA"},
    "board":           {"primary": "1C1C1C", "accent": "B8960C", "light_bg": "FAFAFA"},
    # ── Default fallback ──────────────────────────────────────────────────────
    "default":         {"primary": "2D3142", "accent": "A0522D", "light_bg": "F7F7F9"},
}

# Callout label system — descriptive labels (not fixed badge types)
# The badge displays the callout's own label (e.g. "Key Highlight", "Risk Alert").
# Examples of good descriptive labels:
#   "Key Highlight"    "Market Inflection"   "Critical Finding"   "Risk Alert"
#   "Strategic Vision" "Data Insight"        "Growth Driver"      "Action Required"
#   "Market Scale"     "Closing Insight"     "Verdict"            "Recommendation"
# Legacy short labels (KEY, STAT, INSIGHT, GOAL, NOTE) are still accepted but
# the LLM prompts now request descriptive labels for better document quality.
CALLOUT_LABELS = {
    # Legacy short labels (still accepted)
    "KEY", "STAT", "GOAL", "INSIGHT", "RISK", "NOTE", "VERDICT", "RECOMMENDATION",
    # New descriptive labels (preferred)
    "KEY HIGHLIGHT", "MARKET INFLECTION", "CRITICAL FINDING", "RISK ALERT",
    "STRATEGIC VISION", "DATA INSIGHT", "GROWTH DRIVER", "ACTION REQUIRED",
    "MARKET SCALE", "CLOSING INSIGHT",
}


def get_palette(domain: str | None) -> dict:
    """
    Return the colour palette dict for a given domain keyword.
    Falls back to 'default' if domain is None or unrecognised.
    Accepts compound strings like 'AI / Machine Learning' — tries each word.
    """
    if not domain:
        return DOMAIN_PALETTES["default"]
    key = domain.lower().strip()
    if key in DOMAIN_PALETTES:
        return DOMAIN_PALETTES[key]
    # Try individual words in compound domain strings
    for word in re.split(r"[\s/,_-]+", key):
        if word in DOMAIN_PALETTES:
            return DOMAIN_PALETTES[word]
    return DOMAIN_PALETTES["default"]


# ── DOCX ─────────────────────────────────────────────────────────────────────
#
# Premium executive-quality Word document generator.
# Produces McKinsey / BCG / Claude-style reports with:
#   - Branded cover page (topic-specific colour palette)
#   - Header: org name | document title rule line on every page
#   - Footer: "Confidential — Internal Use Only" + auto page numbers
#   - H1 full-width coloured band headings with thick accent rules
#   - H2 underlined subheadings
#   - Optional H3 subheadings and key-takeaway highlight boxes
#   - Color-coded data tables with alternating row shading
#   - Callout boxes with label badges (KEY / STAT / GOAL / INSIGHT)
#   - Bullet lists with branded arrow markers
#   - Body text: Calibri 11pt, 1.15 line spacing, justified
#
# Libraries used:
#   python-docx  (import docx)  — .docx generation
#   reportlab    (import reportlab) — .pdf generation  (see generate_pdf below)

# ── Fixed body/secondary colours (same across all palettes) ──────────────────
_DOCX_BODY      = "2C2C2C"   # body text (warm near-black, matches Claude)
_DOCX_H1_TX     = "FFFFFF"   # H1 band text (always white)
_DOCX_SECONDARY = "555555"   # secondary / date text (darker for better contrast)
_DOCX_GOLD      = "E8B84B"   # gold accent — callout borders, cover divider

# ── Document footer / confidentiality text ────────────────────────────────────
# DOC_FOOTER_TEXT — text shown in the footer of every generated document.
# OSS default: "Confidential" (generic, no org-specific marking).
# Enterprise: set to e.g. "Acme Corp — Confidential" or "Acme Corp — Internal Only".
_DOC_FOOTER_TEXT = os.getenv("DOC_FOOTER_TEXT", "Confidential")


def generate_docx(title: str, sections: list, domain: str | None = None) -> bytes:
    """
    Generate a premium branded executive Word document (.docx).

    Features:
      - Branded cover page with topic-specific colour palette (domain-driven)
      - Header: "AiNxt  |  <document title>" rule line on every page
      - Footer: "Confidential — Internal Use Only  |  Page N" on every page
      - H1: full-width primary-colour band with thick accent rule beneath
      - H2: primary-colour bold text, 16pt, accent underline rule
      - H3: accent-colour italic subheadings
      - Color-coded data tables: header in primary colour (white text),
        alternating light_bg / white rows, clean grid borders
      - Callout boxes: two-column layout — label badge (accent fill, white text)
        + insight text (light_bg background)
      - Bullet lists with branded → arrow markers
      - Body: Calibri 11pt, 1.15 line spacing, justified

    Sections support the extended schema:
      {
        "heading":      str,   # required
        "subheading":   str,   # optional H3
        "key_takeaway": str,   # optional shaded insight box
        "callout":      {"label": str, "text": str},  # badge callout
        "content":      str,   # body paragraphs (\\n\\n separated)
        "bullets":      list,  # bullet points
        "table":        {"headers": [...], "rows": [[...], ...]},
        "level":        int,   # 1 = H1, 2 = H2 (default)
      }

    Library: python-docx
    """
    logger.info(f"[doc_generator] generate_docx START | title={title!r} sections={len(sections)} domain={domain!r}")
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm, Inches, Mm
        from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        import datetime
        logger.info("[doc_generator] python-docx imported OK")
    except ImportError:
        logger.error("[doc_generator] python-docx NOT installed")
        raise RuntimeError("python-docx not installed — run: pip install python-docx")

    # ── Resolve domain palette ────────────────────────────────────
    pal = get_palette(domain)
    PAL_PRIMARY  = pal["primary"]   # hex, no #
    PAL_ACCENT   = pal["accent"]
    PAL_LIGHT_BG = pal["light_bg"]

    # ── Colour helpers ────────────────────────────────────────────
    def _rgb(hex6: str) -> RGBColor:
        h = hex6.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    C_PRIMARY  = _rgb(PAL_PRIMARY)
    C_ACCENT   = _rgb(PAL_ACCENT)
    C_LIGHT_BG = _rgb(PAL_LIGHT_BG)
    C_BODY     = _rgb(_DOCX_BODY)
    C_H1TX     = _rgb(_DOCX_H1_TX)
    C_SECONDARY = _rgb(_DOCX_SECONDARY)

    def _set_color(run, rgb: RGBColor):
        run.font.color.rgb = rgb

    def _set_line_spacing(para, spacing: float = 1.42):
        pf = para.paragraph_format
        pf.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
        pf.line_spacing = Pt(spacing * 12)

    def _shade_para(para, hex6: str):
        try:
            pPr = para._p.get_or_add_pPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), hex6.lstrip("#"))
            pPr.append(shd)
        except Exception:
            pass

    def _shade_cell(cell, hex6: str):
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), hex6.lstrip("#"))
            tcPr.append(shd)
        except Exception:
            pass

    def _bottom_border(para, hex6: str, sz: str = "12"):
        try:
            pPr = para._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bot = OxmlElement("w:bottom")
            bot.set(qn("w:val"), "single")
            bot.set(qn("w:sz"), sz)
            bot.set(qn("w:space"), "4")
            bot.set(qn("w:color"), hex6.lstrip("#"))
            pBdr.append(bot)
            pPr.append(pBdr)
        except Exception:
            pass

    def _top_border(para, hex6: str, sz: str = "12"):
        try:
            pPr = para._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            top = OxmlElement("w:top")
            top.set(qn("w:val"), "single")
            top.set(qn("w:sz"), sz)
            top.set(qn("w:space"), "4")
            top.set(qn("w:color"), hex6.lstrip("#"))
            pBdr.append(top)
            pPr.append(pBdr)
        except Exception:
            pass

    def _left_border(para, hex6: str, sz: str = "24"):
        try:
            pPr = para._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            left = OxmlElement("w:left")
            left.set(qn("w:val"), "single")
            left.set(qn("w:sz"), sz)
            left.set(qn("w:space"), "6")
            left.set(qn("w:color"), hex6.lstrip("#"))
            pBdr.append(left)
            pPr.append(pBdr)
        except Exception:
            pass

    # ── Header / Footer helpers ───────────────────────────────────

    def _add_header_footer(doc_obj, org_name: str = "AiNxt"):
        """
        Add 3-zone branded header and confidential footer to all sections.

        Header layout (single line):
          LEFT:   [Doc Title] (truncated)
          CENTER: Confidential | © YYYY
          RIGHT:  Page N  (auto field)
        Footer:
          CENTER: Confidential — Internal Use Only  |  Page N
        """
        import datetime as _dt_docx
        _year = _dt_docx.date.today().year

        for sec in doc_obj.sections:
            sec.different_first_page_header_footer = True  # no header on cover

            # ── Header ──────────────────────────────────────────
            hdr = sec.header
            hdr.is_linked_to_previous = False
            for p in hdr.paragraphs:
                p.clear()
            if not hdr.paragraphs:
                hdr.add_paragraph()
            hp = hdr.paragraphs[0]
            # Use a 3-column tab-stop layout: left | center | right
            hp.alignment = WD_ALIGN_PARAGRAPH.LEFT
            hp.paragraph_format.space_before = Pt(0)
            hp.paragraph_format.space_after  = Pt(2)
            _bottom_border(hp, PAL_ACCENT, "4")   # thin 0.5pt border (Claude style)

            # LEFT zone — document title (truncated)
            short_title = title[:60] + ("\u2026" if len(title) > 60 else "")
            title_run = hp.add_run(short_title)
            title_run.font.name = "Arial"
            title_run.font.size = Pt(9)
            title_run.font.bold = False
            title_run.font.italic = True
            _set_color(title_run, C_ACCENT)

            # CENTER zone — tab + confidentiality text
            tab_run = hp.add_run("\t")
            tab_run.font.size = Pt(9)
            conf_run = hp.add_run(f"{_DOC_FOOTER_TEXT}  |  \u00a9 {_year}")
            conf_run.font.name   = "Arial"
            conf_run.font.size   = Pt(8)
            conf_run.font.italic = True
            _set_color(conf_run, _rgb(_DOCX_SECONDARY))

            # RIGHT zone — tab + page number field
            tab_run2 = hp.add_run("\t")
            tab_run2.font.size = Pt(9)
            pg_label = hp.add_run("Page ")
            pg_label.font.name = "Arial"
            pg_label.font.size = Pt(9)
            _set_color(pg_label, _rgb(_DOCX_SECONDARY))
            try:
                fldChar1 = OxmlElement("w:fldChar")
                fldChar1.set(qn("w:fldCharType"), "begin")
                instrText = OxmlElement("w:instrText")
                instrText.text = "PAGE"
                fldChar2 = OxmlElement("w:fldChar")
                fldChar2.set(qn("w:fldCharType"), "end")
                pg_num_run = hp.add_run()
                pg_num_run.font.name = "Arial"
                pg_num_run.font.size = Pt(9)
                _set_color(pg_num_run, _rgb(_DOCX_SECONDARY))
                pg_num_run._r.append(fldChar1)
                pg_num_run._r.append(instrText)
                pg_num_run._r.append(fldChar2)
            except Exception:
                pass

            # Set tab stops: center at ~half page, right at content width
            try:
                from docx.oxml.ns import qn as _qn2
                pPr = hp._p.get_or_add_pPr()
                tabs_elem = OxmlElement("w:tabs")
                # Center tab
                ct = OxmlElement("w:tab")
                ct.set(_qn2("w:val"), "center")
                ct.set(_qn2("w:pos"), "4819")   # ~half of A4 content width in twips (170mm/2)
                tabs_elem.append(ct)
                # Right tab
                rt = OxmlElement("w:tab")
                rt.set(_qn2("w:val"), "right")
                rt.set(_qn2("w:pos"), "9638")   # ~full A4 content width in twips (170mm)
                tabs_elem.append(rt)
                pPr.append(tabs_elem)
            except Exception:
                pass

            # ── Footer ──────────────────────────────────────────
            ftr = sec.footer
            ftr.is_linked_to_previous = False
            for p in ftr.paragraphs:
                p.clear()
            if not ftr.paragraphs:
                ftr.add_paragraph()
            fp = ftr.paragraphs[0]
            fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            fp.paragraph_format.space_before = Pt(2)
            fp.paragraph_format.space_after  = Pt(0)
            _top_border(fp, PAL_ACCENT, "6")

            conf_run = fp.add_run(f"{_DOC_FOOTER_TEXT}  |  Page ")
            conf_run.font.name   = "Arial"
            conf_run.font.size   = Pt(8)
            conf_run.font.italic = True
            _set_color(conf_run, _rgb(_DOCX_SECONDARY))

            # Auto page number field
            try:
                fldChar1 = OxmlElement("w:fldChar")
                fldChar1.set(qn("w:fldCharType"), "begin")
                instrText = OxmlElement("w:instrText")
                instrText.text = "PAGE"
                fldChar2 = OxmlElement("w:fldChar")
                fldChar2.set(qn("w:fldCharType"), "end")
                page_run = fp.add_run()
                page_run.font.name = "Arial"
                page_run.font.size = Pt(8)
                _set_color(page_run, _rgb(_DOCX_SECONDARY))
                page_run._r.append(fldChar1)
                page_run._r.append(instrText)
                page_run._r.append(fldChar2)
            except Exception:
                pass

    # ── Reusable paragraph builders ───────────────────────────────

    def _add_body_para(doc, text: str, indent_cm: float = 0.0,
                       space_before: float = 3, space_after: float = 4,
                       align=None, italic: bool = False):
        p = doc.add_paragraph()
        p.alignment = align or WD_ALIGN_PARAGRAPH.JUSTIFY
        p.paragraph_format.space_before = Pt(space_before)
        p.paragraph_format.space_after  = Pt(space_after)
        if indent_cm:
            p.paragraph_format.left_indent = Cm(indent_cm)
        _set_line_spacing(p, 1.42)   # Claude-matching generous line spacing
        run = p.add_run(text)
        run.font.name   = "Arial"
        run.font.size   = Pt(11)
        run.font.italic = italic
        _set_color(run, C_BODY)
        return p

    def _add_h1(doc, text: str) -> None:
        """H1: primary-colour bold text 17pt, blue bottom border rule (Claude style)."""
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(20)
        p.paragraph_format.space_after  = Pt(8)
        _bottom_border(p, PAL_PRIMARY, "6")   # thin blue underline rule
        run = p.add_run(text)
        run.font.name = "Arial"
        run.font.size = Pt(17)
        run.font.bold = True
        _set_color(run, C_PRIMARY)

    def _add_h2(doc, text: str) -> None:
        """H2: accent-colour bold text, 13pt, no italic (Claude style)."""
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(14)
        p.paragraph_format.space_after  = Pt(5)
        run = p.add_run(text)
        run.font.name   = "Arial"
        run.font.size   = Pt(13)
        run.font.bold   = True
        run.font.italic = False   # Claude uses bold-only, no italic
        _set_color(run, C_ACCENT)

    def _add_h3(doc, text: str) -> None:
        """H3: primary-colour bold, 12pt (Claude style — clean, no italic)."""
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after  = Pt(4)
        run = p.add_run(text)
        run.font.name   = "Arial"
        run.font.size   = Pt(12)
        run.font.bold   = True
        run.font.italic = False
        _set_color(run, C_PRIMARY)

    def _add_bullet(doc, text: str) -> None:
        """Bullet with • marker in accent colour (Arial, Claude style)."""
        p = doc.add_paragraph()
        p.paragraph_format.left_indent       = Cm(1.0)
        p.paragraph_format.first_line_indent = Cm(-0.6)
        p.paragraph_format.space_before      = Pt(3)
        p.paragraph_format.space_after       = Pt(3)
        _set_line_spacing(p, 1.42)
        bullet_run = p.add_run("\u2022  ")   # • bullet
        bullet_run.font.name = "Arial"
        bullet_run.font.size = Pt(11)
        bullet_run.font.bold = True
        _set_color(bullet_run, C_ACCENT)
        text_run = p.add_run(str(text).strip())
        text_run.font.name = "Arial"
        text_run.font.size = Pt(11)
        _set_color(text_run, C_BODY)

    def _add_callout_box(doc, label: str, text: str) -> None:
        """
        Claude-style single-cell callout box:
          - Thick gold left border (sz=18, #E8B84B)
          - Warm cream background (#FFF8E7)
          - Top/bottom/right borders: none (invisible)
          - Bold label inline + italic content text
        """
        spacer = doc.add_paragraph()
        spacer.paragraph_format.space_before = Pt(8)
        spacer.paragraph_format.space_after  = Pt(0)

        # Single-cell table — full text width
        tbl = doc.add_table(rows=1, cols=1)
        tbl.style = "Table Grid"

        # Set explicit full-width (A4 content width at 1-inch margins)
        try:
            from docx.oxml.ns import qn as _qn_tbl
            tblPr = tbl._tbl.get_or_add_tblPr()
            tblW = OxmlElement("w:tblW")
            tblW.set(_qn_tbl("w:w"), "9638")
            tblW.set(_qn_tbl("w:type"), "dxa")
            tblPr.append(tblW)
        except Exception:
            pass

        cell = tbl.rows[0].cells[0]

        # Warm cream background (#FFF8E7)
        _shade_cell(cell, "FFF8E7")

        # Set borders: thick gold left, none on top/bottom/right
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            tcBdr = OxmlElement("w:tcBdr")
            # Left: thick gold (sz=18 = 2.25pt)
            left_bdr = OxmlElement("w:left")
            left_bdr.set(qn("w:val"), "single")
            left_bdr.set(qn("w:sz"), "18")
            left_bdr.set(qn("w:space"), "0")
            left_bdr.set(qn("w:color"), _DOCX_GOLD)
            tcBdr.append(left_bdr)
            # Top: none
            top_bdr = OxmlElement("w:top")
            top_bdr.set(qn("w:val"), "none")
            top_bdr.set(qn("w:sz"), "0")
            top_bdr.set(qn("w:space"), "0")
            top_bdr.set(qn("w:color"), "auto")
            tcBdr.append(top_bdr)
            # Bottom: none
            bot_bdr = OxmlElement("w:bottom")
            bot_bdr.set(qn("w:val"), "none")
            bot_bdr.set(qn("w:sz"), "0")
            bot_bdr.set(qn("w:space"), "0")
            bot_bdr.set(qn("w:color"), "auto")
            tcBdr.append(bot_bdr)
            # Right: none
            right_bdr = OxmlElement("w:right")
            right_bdr.set(qn("w:val"), "none")
            right_bdr.set(qn("w:sz"), "0")
            right_bdr.set(qn("w:space"), "0")
            right_bdr.set(qn("w:color"), "auto")
            tcBdr.append(right_bdr)
            tcPr.append(tcBdr)
            # Cell padding: generous left indent
            tcMar = OxmlElement("w:tcMar")
            for side, val in [("left", "240"), ("right", "120"), ("top", "80"), ("bottom", "80")]:
                m = OxmlElement(f"w:{side}")
                m.set(qn("w:w"), val)
                m.set(qn("w:type"), "dxa")
                tcMar.append(m)
            tcPr.append(tcMar)
        except Exception:
            pass

        cp = cell.paragraphs[0]
        cp.alignment = WD_ALIGN_PARAGRAPH.LEFT
        cp.paragraph_format.space_before = Pt(4)
        cp.paragraph_format.space_after  = Pt(4)
        _set_line_spacing(cp, 1.42)

        # Bold label inline (primary colour)
        label_run = cp.add_run(f"{label.title()}:  ")
        label_run.font.name = "Arial"
        label_run.font.size = Pt(11)
        label_run.font.bold = True
        _set_color(label_run, C_PRIMARY)

        # Italic content text (body colour)
        text_run = cp.add_run(str(text).strip())
        text_run.font.name   = "Arial"
        text_run.font.size   = Pt(10)
        text_run.font.italic = True
        _set_color(text_run, _rgb(_DOCX_BODY))

        spacer2 = doc.add_paragraph()
        spacer2.paragraph_format.space_before = Pt(0)
        spacer2.paragraph_format.space_after  = Pt(8)

    def _add_key_takeaway(doc, text: str) -> None:
        """Shaded insight box with gold left border (Claude style)."""
        spacer = doc.add_paragraph()
        spacer.paragraph_format.space_before = Pt(4)
        spacer.paragraph_format.space_after  = Pt(0)

        p = doc.add_paragraph()
        p.paragraph_format.left_indent  = Cm(0.5)
        p.paragraph_format.right_indent = Cm(0.5)
        p.paragraph_format.space_before = Pt(6)
        p.paragraph_format.space_after  = Pt(6)
        _shade_para(p, "FFF8E7")          # warm cream background
        _left_border(p, _DOCX_GOLD, "18") # gold left border sz=18
        _set_line_spacing(p, 1.15)

        label_run = p.add_run("Key Insight:  ")
        label_run.font.name = "Arial"
        label_run.font.size = Pt(11)
        label_run.font.bold = True
        _set_color(label_run, C_PRIMARY)

        text_run = p.add_run(str(text).strip())
        text_run.font.name   = "Arial"
        text_run.font.size   = Pt(10)
        text_run.font.italic = True
        _set_color(text_run, _rgb(_DOCX_BODY))

        spacer2 = doc.add_paragraph()
        spacer2.paragraph_format.space_before = Pt(0)
        spacer2.paragraph_format.space_after  = Pt(4)

    def _add_divider(doc) -> None:
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.space_after  = Pt(4)
        _bottom_border(p, PAL_ACCENT, "6")

    def _add_table(doc, table: dict) -> None:
        """
        Claude-quality data table:
          - Full text width (9360 twips)
          - Header row: primary colour background, white bold Arial text
          - True alternating rows: light_bg (#F0F6FB tint) / white
          - Explicit light-grey cell borders (CCCCCC)
          - tblHeader flag for proper page-break repeat
          - Generous cell padding (4pt top/bottom, 6pt left/right)
        """
        headers   = table.get("headers") or []
        rows      = table.get("rows") or []
        if not headers and not rows:
            return
        col_count = max(len(headers), max((len(r) for r in rows), default=0))
        if col_count == 0:
            return

        total_rows = (1 if headers else 0) + len(rows)
        tbl = doc.add_table(rows=total_rows, cols=col_count)
        tbl.style = "Table Grid"

        # Set full text width explicitly (A4 content width at 1-inch margins)
        try:
            from docx.oxml.ns import qn as _qn_t
            tblPr = tbl._tbl.get_or_add_tblPr()
            tblW = OxmlElement("w:tblW")
            tblW.set(_qn_t("w:w"), "9638")
            tblW.set(_qn_t("w:type"), "dxa")
            tblPr.append(tblW)
        except Exception:
            pass

        def _set_cell_borders(cell, is_header: bool = False):
            """Apply explicit borders to a cell."""
            try:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                tcBdr = OxmlElement("w:tcBdr")
                border_color = PAL_ACCENT if is_header else "CCCCCC"
                for side in ["top", "left", "bottom", "right"]:
                    b = OxmlElement(f"w:{side}")
                    b.set(qn("w:val"), "single")
                    b.set(qn("w:sz"), "4")
                    b.set(qn("w:space"), "0")
                    b.set(qn("w:color"), border_color)
                    tcBdr.append(b)
                tcPr.append(tcBdr)
                # Cell padding
                tcMar = OxmlElement("w:tcMar")
                for side, val in [("left", "120"), ("right", "120"), ("top", "80"), ("bottom", "80")]:
                    m = OxmlElement(f"w:{side}")
                    m.set(qn("w:w"), val)
                    m.set(qn("w:type"), "dxa")
                    tcMar.append(m)
                tcPr.append(tcMar)
            except Exception:
                pass

        row_idx = 0
        if headers:
            hdr_row = tbl.rows[row_idx]
            # Add tblHeader flag for repeat on page break
            try:
                trPr = hdr_row._tr.get_or_add_trPr()
                tblHeader = OxmlElement("w:tblHeader")
                trPr.append(tblHeader)
            except Exception:
                pass
            for ci, hdr in enumerate(headers[:col_count]):
                cell = hdr_row.cells[ci]
                _shade_cell(cell, PAL_PRIMARY)
                _set_cell_borders(cell, is_header=True)
                for p in cell.paragraphs:
                    p.clear()
                p = cell.paragraphs[0]
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.paragraph_format.space_before = Pt(4)
                p.paragraph_format.space_after  = Pt(4)
                run = p.add_run(str(hdr).strip())
                run.font.name = "Arial"
                run.font.size = Pt(10)
                run.font.bold = True
                _set_color(run, _rgb(_DOCX_H1_TX))
            row_idx += 1

        for ri, data_row in enumerate(rows):
            if row_idx >= total_rows:
                break
            tbl_row = tbl.rows[row_idx]
            # True alternating: odd rows light blue tint, even rows white
            row_bg = "F0F6FB" if ri % 2 == 0 else "FFFFFF"
            for ci, val in enumerate(data_row[:col_count]):
                cell = tbl_row.cells[ci]
                _shade_cell(cell, row_bg)
                _set_cell_borders(cell, is_header=False)
                for p in cell.paragraphs:
                    p.clear()
                p = cell.paragraphs[0]
                p.paragraph_format.space_before = Pt(4)
                p.paragraph_format.space_after  = Pt(4)
                run = p.add_run(str(val).strip())
                run.font.name = "Arial"
                run.font.size = Pt(10)
                _set_color(run, C_BODY)
            row_idx += 1

        sp = doc.add_paragraph()
        sp.paragraph_format.space_before = Pt(6)
        sp.paragraph_format.space_after  = Pt(6)

    # ── Document setup ────────────────────────────────────────────
    doc = Document()

    # A4 page (210×297mm), 1-inch (2.54cm) margins for clean layout
    MARGIN = Cm(2.54)
    for sec in doc.sections:
        sec.page_width    = Mm(210)
        sec.page_height   = Mm(297)
        sec.top_margin    = MARGIN
        sec.bottom_margin = MARGIN
        sec.left_margin   = MARGIN
        sec.right_margin  = MARGIN

    # ── COVER PAGE — Clean Minimal Design ────────────────────────
    #
    # Layout (white page, no full-bleed band):
    #   1. Large push-down spacer (~100pt)
    #   2. Document title — primary colour, 26pt bold, centred
    #   3. Gold divider line (#E8B84B) beneath title
    #   4. Date — secondary colour, 11pt
    #   5. Confidential note
    #
    # ── Push-down spacer ──────────────────────────────────────────
    push_down = doc.add_paragraph()
    push_down.paragraph_format.space_before = Pt(100)
    push_down.paragraph_format.space_after  = Pt(0)
    push_down.add_run("").font.size = Pt(1)

    # Title (primary colour, 26pt bold, centred, with gold bottom border)
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_para.paragraph_format.space_before = Pt(0)
    title_para.paragraph_format.space_after  = Pt(0)
    _bottom_border(title_para, _DOCX_GOLD, "8")   # gold divider beneath title
    tr = title_para.add_run(title)
    tr.font.name = "Arial"
    tr.font.size = Pt(26)
    tr.font.bold = True
    _set_color(tr, C_PRIMARY)

    # Spacer
    sp1 = doc.add_paragraph()
    sp1.paragraph_format.space_before = Pt(42)
    sp1.paragraph_format.space_after  = Pt(0)

    # Date line
    date_para = doc.add_paragraph()
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    date_para.paragraph_format.space_before = Pt(0)
    date_para.paragraph_format.space_after  = Pt(14)
    today = datetime.date.today().strftime("%B %Y")
    dr = date_para.add_run(today)
    dr.font.name = "Arial"
    dr.font.size = Pt(11)
    _set_color(dr, C_SECONDARY)

    # Confidential note
    conf_para = doc.add_paragraph()
    conf_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    conf_para.paragraph_format.space_before = Pt(8)
    conf_para.paragraph_format.space_after  = Pt(0)
    conf_r = conf_para.add_run(_DOC_FOOTER_TEXT)
    conf_r.font.name   = "Arial"
    conf_r.font.size   = Pt(9)
    conf_r.font.italic = True
    _set_color(conf_r, C_SECONDARY)

    # ── Add header/footer to all sections ─────────────────────────
    _add_header_footer(doc)

    # ── BODY SECTIONS — no TOC, content starts immediately ────────
    doc.add_page_break()

    # Section numbering counters
    _h1_counter = 0
    _h2_counter = 0

    for sec_idx, sec in enumerate(sections):
        h              = (sec.get("heading")        or "").strip()
        subheading     = (sec.get("subheading")     or "").strip()
        key_takeaway   = (sec.get("key_takeaway")   or "").strip()
        callout        = sec.get("callout")  # {"label": str, "text": str}
        verdict        = (sec.get("verdict")        or "").strip()
        recommendation = (sec.get("recommendation") or "").strip()
        key_insight    = (sec.get("key_insight")    or "").strip()
        content        = (sec.get("content")        or "").strip()
        bullets        = sec.get("bullets") or []
        table          = sec.get("table")
        level          = int(sec.get("level") or 2)

        # Heading with section numbering.
        # Strip any leading number/bullet already in the heading text (e.g. "1)", "1.", "1 ")
        # to avoid double-numbering like "1.  1) Executive Summary".
        def _strip_leading_number(text: str) -> str:
            return re.sub(r"^\s*\d+[\.\)]\s*", "", text).strip()

        if h:
            if level == 1:
                _h1_counter += 1
                _h2_counter  = 0
                _add_h1(doc, f"{_h1_counter}.  {_strip_leading_number(h)}")
            else:
                _h2_counter += 1
                _add_h2(doc, f"{_h1_counter}.{_h2_counter}  {_strip_leading_number(h)}")

        if subheading:
            _add_h3(doc, subheading)

        if content:
            for para_text in [p.strip() for p in content.split("\n\n") if p.strip()]:
                _add_body_para(doc, para_text)

        # Callout box — use the callout's own descriptive label
        if callout and isinstance(callout, dict):
            _add_callout_box(doc, callout.get("label", "Key Highlight"), callout.get("text", ""))
        elif key_takeaway:
            _add_key_takeaway(doc, key_takeaway)

        # Legacy inline callouts
        if key_insight:
            _add_callout_box(doc, "Key Insight", key_insight)
        if verdict:
            _add_callout_box(doc, "Verdict", verdict)
        if recommendation:
            _add_callout_box(doc, "Strategic Goal", recommendation)

        for b in bullets:
            if b and str(b).strip():
                _add_bullet(doc, str(b))

        if table and isinstance(table, dict):
            _add_table(doc, table)

        if level == 1 and sec_idx < len(sections) - 1:
            _add_divider(doc)

    buf = io.BytesIO()
    doc.save(buf)
    result_bytes = buf.getvalue()
    logger.info(f"[doc_generator] generate_docx DONE | size={len(result_bytes):,} bytes domain={domain!r} palette={PAL_PRIMARY}")
    return result_bytes


# ── PPTX ─────────────────────────────────────────────────────────────────────
# All helpers are defined before generate_pptx() and are module-private.

try:
    from pptx import Presentation as _Presentation
    from pptx.util import Inches as _Inches, Pt as _Pt, Emu as _Emu
    from pptx.dml.color import RGBColor as _RGB
    from pptx.enum.text import PP_ALIGN as _ALIGN, MSO_ANCHOR as _ANCHOR
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree

    # Default PPTX colour palette
    _NAVY   = _RGB(0x00, 0x33, 0x66)
    _ORANGE = _RGB(0xFF, 0x66, 0x00)
    _WHITE  = _RGB(0xFF, 0xFF, 0xFF)
    _LIGHT  = _RGB(0xF5, 0xF7, 0xFA)
    _GRAY   = _RGB(0x33, 0x33, 0x44)
    _GOLD   = _RGB(0xFF, 0xB3, 0x33)
    _STEEL  = _RGB(0x99, 0xBB, 0xDD)
    _DARK   = _RGB(0x0A, 0x1A, 0x2E)

    _PPTX_NS    = "http://schemas.openxmlformats.org/presentationml/2006/main"
    _DRAWINGML  = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # Default gradient stops (used when no theme is passed to a renderer)
    _DEFAULT_DARK_STOPS  = [(0, "060D1A"), (50, "0A1A2E"), (100, "003366")]
    _DEFAULT_LIGHT_STOPS = [(0, "F5F7FA"), (100, "EAF0FF")]
    _DEFAULT_DARK_ANGLE  = 135   # pptx clockwise-from-East
    _DEFAULT_LIGHT_ANGLE = 90    # pptx clockwise-from-East

    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


# ── Low-level primitives ──────────────────────────────────────────────────────

def _set_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _set_gradient_bg(slide, stops, angle_deg=90):
    """
    Set slide background to a linear gradient using python-pptx native API.
    stops : [(pos_pct 0-100, hex_color_no_hash), ...]
    angle_deg : pptx convention (clockwise from East) — 90=top→bottom, 135=diagonal NW→SE
    Falls back to solid fill on any error.
    """
    try:
        fill = slide.background.fill
        fill.gradient()  # native API — generates proper rotWithShape="1"

        # Set angle: pptx gradient_angle is clockwise from East
        fill.gradient_angle = angle_deg

        # Set the two built-in stops (first and last)
        fill.gradient_stops[0].position = stops[0][0] / 100.0
        fill.gradient_stops[0].color.rgb = _RGB(*bytes.fromhex(stops[0][1]))
        fill.gradient_stops[1].position = stops[-1][0] / 100.0
        fill.gradient_stops[1].color.rgb = _RGB(*bytes.fromhex(stops[-1][1]))

        # Inject any intermediate stops via lxml on the already-valid gsLst
        if len(stops) > 2:
            bgPr = fill._xPr
            gsLst = bgPr.find(f'.//{_qn("a:gsLst")}')
            for i, (pos, hex_color) in enumerate(stops[1:-1], start=1):
                gs = _etree.SubElement(gsLst, _qn("a:gs"))
                gs.set("pos", str(int(pos * 1000)))
                srgb = _etree.SubElement(gs, _qn("a:srgbClr"))
                srgb.set("val", hex_color)
                gsLst.remove(gs)
                gsLst.insert(i, gs)

    except Exception as exc:
        logger.warning(f"doc_generator: gradient bg failed: {exc}")
        if stops:
            try:
                _set_bg(slide, _RGB(*bytes.fromhex(stops[0][1])))
            except Exception:
                pass


def _add_rect(slide, left, top, width, height, fill_rgb, line=False, alpha=None):
    sp = slide.shapes.add_shape(1, left, top, width, height)  # 1 = RECTANGLE
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_rgb
    if not line:
        sp.line.fill.background()
    if alpha is not None:
        _set_alpha(sp, alpha)
    return sp


def _set_alpha(shape, alpha_pct: int):
    """Set shape fill opacity. alpha_pct: 0=transparent, 100=opaque."""
    try:
        # shape.fill._fill IS the <a:solidFill> element — look for color directly inside it
        fill_elem = shape.fill._fill
        clr = (fill_elem.find(_qn("a:srgbClr")) or
               fill_elem.find(_qn("a:sysClr")) or
               fill_elem.find(_qn("a:schemeClr")))
        if clr is None:
            return
        for old in clr.findall(_qn("a:alpha")):
            clr.remove(old)
        a_elem = _etree.SubElement(clr, _qn("a:alpha"))
        a_elem.set("val", str(int(alpha_pct * 1000)))
    except Exception:
        pass


def _add_textbox(slide, left, top, width, height, text,
                 font_name="Calibri", font_pt=18, bold=False, italic=False,
                 color=None, align=None, v_anchor=None, wrap=True):
    if color is None:
        color = _WHITE
    if align is None:
        align = _ALIGN.LEFT
    if v_anchor is None:
        v_anchor = _ANCHOR.TOP

    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    tf.vertical_anchor = v_anchor
    tf.margin_left   = _Pt(3)
    tf.margin_right  = _Pt(3)
    tf.margin_top    = _Pt(2)
    tf.margin_bottom = _Pt(2)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    f = run.font
    f.name   = font_name
    f.size   = _Pt(font_pt)
    f.bold   = bold
    f.italic = italic
    f.color.rgb = color
    return txb


def _add_bullets(tf, bullets, font_pt=16, color=None, bullet_char="→", font_name="Calibri"):
    if color is None:
        color = _GRAY
    for i, text in enumerate(bullets):
        if not text or not str(text).strip():
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = _ALIGN.LEFT
        p.space_before = _Pt(7)
        run = p.add_run()
        run.text = f"{bullet_char}  {str(text).strip()}"
        f = run.font
        f.name  = font_name
        f.size  = _Pt(font_pt)
        f.bold  = False
        f.color.rgb = color


def _add_bg_image(slide, img_bytes: bytes, prs):
    """Add image as full-bleed background, z-ordered behind all other shapes."""
    pic = slide.shapes.add_picture(
        io.BytesIO(img_bytes), 0, 0,
        width=prs.slide_width, height=prs.slide_height,
    )
    sp_tree = slide.shapes._spTree
    pic_elem = pic._element
    sp_tree.remove(pic_elem)
    sp_tree.insert(2, pic_elem)  # behind everything
    return pic


def _add_right_panel_image(slide, img_bytes: bytes, x_frac=0.60, prs=None):
    """Add image as right-side panel (x_frac = left edge as fraction of slide width)."""
    w = prs.slide_width
    h = prs.slide_height
    left  = int(w * x_frac)
    width = w - left
    pic = slide.shapes.add_picture(io.BytesIO(img_bytes), left, 0, width=width, height=h)
    sp_tree = slide.shapes._spTree
    pic_elem = pic._element
    sp_tree.remove(pic_elem)
    sp_tree.insert(2, pic_elem)
    return pic


def _mark(shape, name: str):
    """Tag a shape for animation by setting its name."""
    shape.name = name


# ── Geometric accent helpers (used when no image is available) ────────────────

def _add_decorative_circles(slide, prs):
    """Large semi-transparent accent circles for visual depth on dark slides."""
    w, h = prs.slide_width, prs.slide_height
    # Big circle top-right (bleeds off slide edge)
    _add_rect(slide, int(w * 0.72), int(-h * 0.25), int(w * 0.5), int(w * 0.5),
              _ORANGE, alpha=18)
    # Medium circle bottom-left
    _add_rect(slide, int(-w * 0.04), int(h * 0.55), int(w * 0.32), int(w * 0.32),
              _STEEL, alpha=15)
    # Small accent dot top-left
    _add_rect(slide, int(w * 0.02), int(h * 0.05), int(w * 0.06), int(w * 0.06),
              _ORANGE, alpha=25)


def _add_light_grid_pattern(slide, prs):
    """Lightweight accent pattern on light-background slides (diagonal lines only)."""
    w, h = prs.slide_width, prs.slide_height
    # Three subtle diagonal accent lines in bottom-right corner area
    for i in range(3):
        x_start = int(w * (0.78 + i * 0.07))
        _add_rect(slide, x_start, int(h * 0.55), _Inches(0.04), int(h * 0.5),
                  _NAVY, alpha=8)


# ── Slide type renderers ──────────────────────────────────────────────────────

def _render_title(slide, data, prs_title, prs, tc=None):
    tc = tc or {}
    img_bytes = data.get("_image_bytes")

    if img_bytes:
        _add_bg_image(slide, img_bytes, prs)
        # Semi-transparent dark overlay — text layers render above it
        _add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, _DARK, alpha=74)
    else:
        _set_gradient_bg(slide,
                         tc.get("dark_stops", _DEFAULT_DARK_STOPS),
                         tc.get("dark_angle", _DEFAULT_DARK_ANGLE))
        _add_decorative_circles(slide, prs)

    # Orange accent bar (bottom third line)
    _add_rect(slide, 0, _Inches(4.15), prs.slide_width, _Inches(0.1), _ORANGE)

    # Main heading
    heading_tb = _add_textbox(
        slide, _Inches(0.6), _Inches(1.3), _Inches(8.8), _Inches(1.7),
        text=data.get("heading") or prs_title,
        font_pt=46, bold=True, color=_WHITE,
        align=_ALIGN.CENTER, v_anchor=_ANCHOR.MIDDLE,
    )
    _mark(heading_tb, "anim_0")

    # Key message / subtitle
    km = (data.get("key_message") or "").strip()
    if km:
        km_tb = _add_textbox(
            slide, _Inches(0.8), _Inches(3.1), _Inches(8.4), _Inches(0.75),
            text=km, font_pt=22, color=_GOLD,
            align=_ALIGN.CENTER,
        )
        _mark(km_tb, "anim_1")

    # Icon
    icon = (data.get("icon") or "").strip()
    if icon:
        _add_textbox(
            slide, _Inches(8.5), _Inches(0.2), _Inches(1.2), _Inches(0.9),
            text=icon, font_pt=38, color=_WHITE, align=_ALIGN.CENTER,
        )

    # Footer accent bar
    _add_rect(slide, 0, _Inches(5.15), prs.slide_width, _Inches(0.475), _ORANGE, alpha=90)
    _add_textbox(
        slide, 0, _Inches(5.15), prs.slide_width, _Inches(0.475),
        text="AiNxt AI Platform",
        font_pt=12, bold=True, color=_WHITE,
        align=_ALIGN.CENTER, v_anchor=_ANCHOR.MIDDLE,
    )


def _render_agenda(slide, data, prs_title, prs, tc=None):
    tc = tc or {}
    _set_gradient_bg(slide,
                     tc.get("light_stops", _DEFAULT_LIGHT_STOPS),
                     tc.get("light_angle", _DEFAULT_LIGHT_ANGLE))
    _add_light_grid_pattern(slide, prs)

    # Navy header band
    _add_rect(slide, 0, 0, prs.slide_width, _Inches(1.05), _NAVY)
    # Decorative orange sliver at top of header
    _add_rect(slide, 0, 0, prs.slide_width, _Inches(0.07), _ORANGE)

    hdr_tb = _add_textbox(
        slide, _Inches(0.4), 0, _Inches(9.2), _Inches(1.05),
        text=data.get("heading") or "Agenda",
        font_pt=28, bold=True, color=_WHITE,
        v_anchor=_ANCHOR.MIDDLE,
    )
    _mark(hdr_tb, "anim_0")

    # Orange left accent strip
    _add_rect(slide, 0, _Inches(1.05), _Inches(0.14), _Inches(4.575), _ORANGE)

    bullets = data.get("bullets") or []
    y_start = 1.3
    row_h   = 0.68

    for i, item in enumerate(bullets[:6]):
        y = y_start + i * row_h

        # Orange numbered badge
        badge = slide.shapes.add_shape(1,
            _Inches(0.28), _Inches(y),
            _Inches(0.38), _Inches(0.38))
        badge.fill.solid()
        badge.fill.fore_color.rgb = _ORANGE
        badge.line.fill.background()
        btf = badge.text_frame
        btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = _Pt(0)
        btf.vertical_anchor = _ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = _ALIGN.CENTER
        br = bp.add_run()
        br.text = str(i + 1)
        br.font.name  = "Calibri"
        br.font.size  = _Pt(14)
        br.font.bold  = True
        br.font.color.rgb = _WHITE

        item_tb = _add_textbox(
            slide, _Inches(0.82), _Inches(y - 0.04), _Inches(8.8), _Inches(0.5),
            text=str(item).strip(),
            font_pt=20, color=_GRAY, v_anchor=_ANCHOR.MIDDLE,
        )
        _mark(item_tb, f"anim_click_{i:02d}")

    # Bottom accent line
    _add_rect(slide, _Inches(0.28), _Inches(5.2), _Inches(9.5), _Inches(0.05), _ORANGE, alpha=60)


def _render_content(slide, data, prs_title, prs, tc=None):
    tc = tc or {}
    img_bytes = data.get("_image_bytes")

    has_image   = img_bytes is not None
    content_w   = _Inches(5.9) if has_image else _Inches(9.5)
    bullets_w   = _Inches(5.6) if has_image else _Inches(9.1)
    header_w    = prs.slide_width

    _set_gradient_bg(slide,
                     tc.get("light_stops", _DEFAULT_LIGHT_STOPS),
                     tc.get("light_angle", _DEFAULT_LIGHT_ANGLE))

    if has_image:
        _add_right_panel_image(slide, img_bytes, x_frac=0.61, prs=prs)
        # Subtle gradient fade overlay between content and image
        fade = _add_rect(slide, _Inches(5.5), 0, _Inches(0.7), prs.slide_height,
                         _LIGHT, alpha=80)
    else:
        _add_light_grid_pattern(slide, prs)

    # Navy header band (full width)
    _add_rect(slide, 0, 0, header_w, _Inches(1.05), _NAVY)
    _add_rect(slide, 0, 0, header_w, _Inches(0.07), _ORANGE)

    # Icon in header
    icon = (data.get("icon") or "").strip()
    hdr_right = _Inches(8.3) if icon else _Inches(9.5)
    hdr_tb = _add_textbox(
        slide, _Inches(0.35), 0, hdr_right, _Inches(1.05),
        text=data.get("heading") or prs_title,
        font_pt=26, bold=True, color=_WHITE, v_anchor=_ANCHOR.MIDDLE,
    )
    _mark(hdr_tb, "anim_0")
    if icon:
        _add_textbox(
            slide, _Inches(8.75), _Inches(0.12), _Inches(0.9), _Inches(0.8),
            text=icon, font_pt=30, color=_WHITE, align=_ALIGN.CENTER,
        )

    # Orange left accent strip
    _add_rect(slide, 0, _Inches(1.05), _Inches(0.13), _Inches(4.575), _ORANGE)

    # Key message pull-quote
    km = (data.get("key_message") or "").strip()
    bullets_top = _Inches(1.25)
    if km:
        km_tb = _add_textbox(
            slide, _Inches(0.35), _Inches(1.18), content_w, _Inches(0.65),
            text=km, font_pt=17, bold=True, color=_NAVY,
        )
        _mark(km_tb, "anim_1")
        _add_rect(slide, _Inches(0.35), _Inches(1.87), content_w, _Inches(0.04), _ORANGE)
        bullets_top = _Inches(2.02)

    # Bullets
    bullets = data.get("bullets") or []
    if bullets:
        txb = slide.shapes.add_textbox(
            _Inches(0.5), bullets_top, bullets_w, _Inches(3.3))
        tf = txb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        _add_bullets(tf, bullets, font_pt=16, color=_GRAY)
        _mark(txb, "anim_click_00")


def _render_stats(slide, data, prs_title, prs, tc=None):
    tc = tc or {}
    _set_gradient_bg(slide,
                     tc.get("dark_stops", _DEFAULT_DARK_STOPS),
                     tc.get("dark_angle", _DEFAULT_DARK_ANGLE))
    _add_decorative_circles(slide, prs)

    w = prs.slide_width

    # Header
    _add_rect(slide, 0, 0, w, _Inches(0.07), _ORANGE)
    hdr_tb = _add_textbox(
        slide, _Inches(0.4), _Inches(0.15), _Inches(9.2), _Inches(0.85),
        text=data.get("heading") or "Key Metrics",
        font_pt=28, bold=True, color=_WHITE, align=_ALIGN.CENTER,
    )
    _mark(hdr_tb, "anim_0")
    _add_rect(slide, _Inches(1.5), _Inches(1.1), _Inches(7.0), _Inches(0.06), _ORANGE)

    stats = data.get("stats") or []
    n = min(len(stats), 3)
    if n == 0:
        return

    col_w = 10.0 / n
    for i, stat in enumerate(stats[:3]):
        x = i * col_w

        # Subtle circle behind value
        cx = _Inches(x + col_w / 2 - 0.8)
        _add_rect(slide, cx, _Inches(1.2), _Inches(1.6), _Inches(1.6), _NAVY, alpha=60)

        val_tb = _add_textbox(
            slide,
            _Inches(x + 0.12), _Inches(1.3), _Inches(col_w - 0.24), _Inches(1.65),
            text=str(stat.get("value") or ""),
            font_pt=54, bold=True, color=_ORANGE,
            align=_ALIGN.CENTER, v_anchor=_ANCHOR.MIDDLE,
        )
        _mark(val_tb, f"anim_click_{i:02d}")

        lbl_tb = _add_textbox(
            slide,
            _Inches(x + 0.12), _Inches(3.05), _Inches(col_w - 0.24), _Inches(0.65),
            text=str(stat.get("label") or ""),
            font_pt=15, color=_WHITE, align=_ALIGN.CENTER,
        )
        _mark(lbl_tb, f"anim_click_{i:02d}_lbl")

        if i < n - 1:
            _add_rect(slide,
                _Inches(x + col_w - 0.03), _Inches(1.3),
                _Inches(0.06), _Inches(2.2),
                _STEEL, alpha=30)

    # Key message
    km = (data.get("key_message") or "").strip()
    if km:
        _add_textbox(
            slide, _Inches(1.0), _Inches(3.9), _Inches(8.0), _Inches(0.65),
            text=km, font_pt=14, italic=True, color=_STEEL, align=_ALIGN.CENTER,
        )

    # Footer accent
    _add_rect(slide, 0, _Inches(5.15), w, _Inches(0.475), _ORANGE, alpha=30)


def _render_quote(slide, data, prs_title, prs, tc=None):
    tc = tc or {}
    _set_gradient_bg(slide,
                     tc.get("dark_stops", _DEFAULT_DARK_STOPS),
                     tc.get("dark_angle", _DEFAULT_DARK_ANGLE))
    _add_decorative_circles(slide, prs)

    # Large decorative quote mark (background element)
    _add_textbox(
        slide, _Inches(0.3), _Inches(0.3), _Inches(2.0), _Inches(2.0),
        text="❝",  # ❝
        font_pt=96, color=_ORANGE, align=_ALIGN.LEFT,
    )

    # Optional context heading
    heading = (data.get("heading") or "").strip()
    if heading:
        _add_textbox(
            slide, _Inches(0.5), _Inches(0.2), _Inches(9.0), _Inches(0.45),
            text=heading, font_pt=13, color=_STEEL, align=_ALIGN.RIGHT,
        )

    # Quote text
    quote_text = (data.get("quote") or "").strip()
    if quote_text:
        q_tb = _add_textbox(
            slide, _Inches(1.0), _Inches(0.95), _Inches(8.0), _Inches(3.0),
            text=f'"{quote_text}"',
            font_pt=24, italic=True, color=_WHITE,
            align=_ALIGN.CENTER, v_anchor=_ANCHOR.MIDDLE,
        )
        _mark(q_tb, "anim_0")

    # Orange separator
    _add_rect(slide, _Inches(3.0), _Inches(4.05), _Inches(4.0), _Inches(0.07), _ORANGE)

    # Attribution
    attribution = (data.get("attribution") or "").strip()
    if attribution:
        attr_tb = _add_textbox(
            slide, _Inches(2.5), _Inches(4.25), _Inches(5.0), _Inches(0.65),
            text=attribution, font_pt=16, italic=True, color=_GOLD, align=_ALIGN.CENTER,
        )
        _mark(attr_tb, "anim_1")


def _render_two_column(slide, data, prs_title, prs, tc=None):
    tc = tc or {}
    _set_gradient_bg(slide,
                     tc.get("light_stops", _DEFAULT_LIGHT_STOPS),
                     tc.get("light_angle", _DEFAULT_LIGHT_ANGLE))
    _add_light_grid_pattern(slide, prs)

    w = prs.slide_width

    # Header band
    _add_rect(slide, 0, 0, w, _Inches(1.0), _NAVY)
    _add_rect(slide, 0, 0, w, _Inches(0.07), _ORANGE)
    hdr_tb = _add_textbox(
        slide, _Inches(0.3), 0, _Inches(9.4), _Inches(1.0),
        text=data.get("heading") or "Comparison",
        font_pt=26, bold=True, color=_WHITE,
        align=_ALIGN.CENTER, v_anchor=_ANCHOR.MIDDLE,
    )
    _mark(hdr_tb, "anim_0")

    # Center vertical divider
    _add_rect(slide, _Inches(4.95), _Inches(1.1), _Inches(0.1), _Inches(4.2), _ORANGE)

    left_data  = data.get("two_col_left")  or {}
    right_data = data.get("two_col_right") or {}

    for col, x_off in [(left_data, 0.3), (right_data, 5.2)]:
        col_title   = (col.get("title") or "").strip()
        col_bullets = col.get("bullets") or []

        if col_title:
            ct_tb = _add_textbox(
                slide, _Inches(x_off), _Inches(1.15), _Inches(4.4), _Inches(0.5),
                text=col_title, font_pt=18, bold=True, color=_NAVY,
            )
            _mark(ct_tb, f"anim_click_{'00' if x_off < 4 else '01'}")
            _add_rect(slide, _Inches(x_off), _Inches(1.69), _Inches(4.4), _Inches(0.05), _ORANGE)

        if col_bullets:
            txb = slide.shapes.add_textbox(
                _Inches(x_off), _Inches(1.88), _Inches(4.5), _Inches(3.3))
            tf = txb.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            _add_bullets(tf, col_bullets, font_pt=15, color=_GRAY, bullet_char="•")
            _mark(txb, f"anim_click_{'00' if x_off < 4 else '01'}_body")


def _render_closing(slide, data, prs_title, prs, tc=None):
    tc = tc or {}
    img_bytes = data.get("_image_bytes")

    if img_bytes:
        _add_bg_image(slide, img_bytes, prs)
        _add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, _DARK, alpha=72)
    else:
        _set_gradient_bg(slide,
                         tc.get("dark_stops", _DEFAULT_DARK_STOPS),
                         tc.get("dark_angle", _DEFAULT_DARK_ANGLE))
        _add_decorative_circles(slide, prs)

    # Thin orange top bar
    _add_rect(slide, 0, 0, prs.slide_width, _Inches(0.1), _ORANGE)

    # Icon
    icon = (data.get("icon") or "🙏").strip()
    _add_textbox(
        slide, _Inches(4.3), _Inches(0.65), _Inches(1.4), _Inches(1.0),
        text=icon, font_pt=36, color=_WHITE, align=_ALIGN.CENTER,
    )

    # Heading
    heading_tb = _add_textbox(
        slide, _Inches(0.6), _Inches(0.8), _Inches(8.8), _Inches(1.3),
        text=data.get("heading") or "Thank You",
        font_pt=42, bold=True, color=_WHITE,
        align=_ALIGN.CENTER, v_anchor=_ANCHOR.MIDDLE,
    )
    _mark(heading_tb, "anim_0")

    km = (data.get("key_message") or "").strip()
    if km:
        km_tb = _add_textbox(
            slide, _Inches(0.6), _Inches(2.15), _Inches(8.8), _Inches(0.65),
            text=km, font_pt=20, color=_GOLD, align=_ALIGN.CENTER,
        )
        _mark(km_tb, "anim_1")

    # Separator
    _add_rect(slide, _Inches(2.5), _Inches(2.95), _Inches(5.0), _Inches(0.08), _ORANGE)

    bullets = data.get("bullets") or []
    if bullets:
        txb = slide.shapes.add_textbox(
            _Inches(2.5), _Inches(3.15), _Inches(5.0), _Inches(1.8))
        tf = txb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        _add_bullets(tf, bullets, font_pt=16, color=_WHITE, bullet_char="✓")
        _mark(txb, "anim_click_00")

    # Footer accent band
    _add_rect(slide, 0, _Inches(5.15), prs.slide_width, _Inches(0.475), _ORANGE)
    _add_textbox(
        slide, 0, _Inches(5.15), prs.slide_width, _Inches(0.475),
        text="AiNxt AI Platform",
        font_pt=12, bold=True, color=_WHITE,
        align=_ALIGN.CENTER, v_anchor=_ANCHOR.MIDDLE,
    )


_SLIDE_RENDERERS = {
    "title":      _render_title,
    "agenda":     _render_agenda,
    "content":    _render_content,
    "stats":      _render_stats,
    "quote":      _render_quote,
    "two_column": _render_two_column,
    "two_col":    _render_two_column,
    "closing":    _render_closing,
    "thank_you":  _render_closing,
} if _PPTX_AVAILABLE else {}


def _legacy_to_rich(title: str, sections: list) -> list:
    """Convert old flat {heading, content, bullets} schema to new rich slide dicts."""
    rich = [{
        "slide_type": "title",
        "heading": title,
        "key_message": "Generated by AiNxt",
        "icon": "📊",
        "speaker_notes": "",
    }]
    for sec in sections:
        rich.append({
            "slide_type":  "content",
            "heading":     (sec.get("heading") or title).strip(),
            "key_message": (sec.get("content") or "").strip()[:120],
            "bullets":     sec.get("bullets") or [],
            "icon":        "",
            "speaker_notes": "",
        })
    rich.append({
        "slide_type":  "closing",
        "heading":     "Thank You",
        "key_message": "Questions & Discussion",
        "bullets":     [],
        "icon":        "🙏",
        "speaker_notes": "",
    })
    return rich


# ── Transitions (injected as XML into slide._element) ────────────────────────

_TRANSITION_MAP = {
    "title":      ('800',  '<p:fade/>'),
    "agenda":     ('500',  '<p:push dir="r"/>'),
    "content":    ('500',  '<p:push dir="r"/>'),
    "stats":      ('600',  '<p:fade/>'),
    "quote":      ('700',  '<p:fade/>'),
    "two_column": ('500',  '<p:push dir="r"/>'),
    "two_col":    ('500',  '<p:push dir="r"/>'),
    "closing":    ('900',  '<p:fade/>'),
}

def _apply_transition(slide, slide_type: str):
    dur, inner = _TRANSITION_MAP.get(slide_type, ('500', '<p:push dir="r"/>'))
    xml_str = (
        f'<p:transition xmlns:p="{_PPTX_NS}" dur="{dur}" '
        f'xmlns:a="{_DRAWINGML}">{inner}</p:transition>'
    )
    try:
        trans_elem = _etree.fromstring(xml_str)
        slide._element.append(trans_elem)
    except Exception as exc:
        logger.warning(f"doc_generator: transition inject failed: {exc}")


# ── Animations (injected as <p:timing> XML) ───────────────────────────────────

def _apply_animations(slide):
    """
    Build <p:timing> XML for all shapes tagged anim_*.
    - anim_0, anim_1 …   : auto-play fade-in (staggered by 400ms each)
    - anim_click_NN …    : on-click entrance (Fade, presetID=10)
    """
    try:
        tagged = [(s.name, s.shape_id) for s in slide.shapes
                  if s.name and s.name.startswith("anim_")]
        if not tagged:
            return
        tagged.sort(key=lambda x: x[0])

        auto_shapes  = [(n, sid) for n, sid in tagged if not n.startswith("anim_click_")]
        click_shapes = [(n, sid) for n, sid in tagged if n.startswith("anim_click_")]

        timing_xml = _build_timing_xml(auto_shapes, click_shapes)
        slide._element.append(_etree.fromstring(timing_xml))
    except Exception as exc:
        logger.warning(f"doc_generator: animation inject failed: {exc}")


def _build_timing_xml(auto_shapes: list, click_shapes: list) -> str:
    """
    Generate a complete <p:timing> element for fade-in animations.

    auto_shapes:  list of (name, shape_id) — animate automatically with delay
    click_shapes: list of (name, shape_id) — animate on click
    """
    P   = _PPTX_NS
    A   = _DRAWINGML

    # Start unique IDs from 1
    uid = [1]
    def next_id():
        v = uid[0]; uid[0] += 1; return v

    def fade_par(shape_id: int, delay_ms: int, grp_id: int, node_type: str = "clickEffect") -> str:
        cTn_id   = next_id()
        set_cTn  = next_id()
        eff_cTn  = next_id()
        delay_str = str(delay_ms) if delay_ms > 0 else "0"
        return f"""
        <p:par xmlns:p="{P}" xmlns:a="{A}">
          <p:cTn id="{cTn_id}" presetID="10" presetClass="entr" presetSubtype="0"
                 fill="hold" grpId="{grp_id}" nodeType="{node_type}">
            <p:stCondLst><p:cond delay="{delay_str}"/></p:stCondLst>
            <p:childTnLst>
              <p:set>
                <p:cBhvr>
                  <p:cTn id="{set_cTn}" dur="1" fill="hold"/>
                  <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr>
                <p:to><p:strVal val="visible"/></p:to>
              </p:set>
              <p:animEffect transition="in" filter="fade">
                <p:cBhvr>
                  <p:cTn id="{eff_cTn}" dur="600"/>
                  <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                </p:cBhvr>
              </p:animEffect>
            </p:childTnLst>
          </p:cTn>
        </p:par>"""

    # Build auto-play sequence paragraphs
    auto_pars = ""
    for idx, (name, sid) in enumerate(auto_shapes):
        delay = idx * 400
        auto_pars += fade_par(sid, delay, idx, node_type="withEffect" if idx > 0 else "afterEffect")

    # Build click-trigger sequence paragraphs
    click_pars = ""
    for idx, (name, sid) in enumerate(click_shapes):
        cTn_outer = next_id()
        inner     = fade_par(sid, 0, len(auto_shapes) + idx, node_type="clickEffect")
        click_pars += f"""
        <p:par xmlns:p="{P}" xmlns:a="{A}">
          <p:cTn id="{cTn_outer}" fill="hold">
            <p:stCondLst><p:cond delay="indefin"/></p:stCondLst>
            <p:childTnLst>{inner}</p:childTnLst>
          </p:cTn>
        </p:par>"""

    root_id   = next_id()
    seq_id    = next_id()
    all_inner = auto_pars + click_pars

    timing = f"""
<p:timing xmlns:p="{P}" xmlns:a="{A}">
  <p:tnLst>
    <p:par>
      <p:cTn id="{root_id}" dur="indefin" restart="whenNotActive" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="{seq_id}" dur="indefin" nodeType="mainSeq">
              <p:childTnLst>{all_inner}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrevClick" delay="0">
                <p:tn><p:prevAc><p:endSync evt="end" delay="0">
                  <p:rtn><p:tnEl><p:cTn id="{seq_id}"/></p:tnEl></p:rtn>
                </p:endSync></p:prevAc></p:tn>
              </p:cond>
            </p:prevCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
  <p:bldLst>
    {"".join(f'<p:bldP spid="{sid}" grpId="{i}" uiExpand="1" build="p"/>' for i, (_, sid) in enumerate(click_shapes))}
  </p:bldLst>
</p:timing>"""
    return timing.strip()


# ── Main PPTX generator ───────────────────────────────────────────────────────

def generate_pptx(title: str, sections: list, use_template: bool = False,
                  theme: str = "dark_executive") -> bytes:
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not installed — run: pip install python-pptx")

    # Fix #37: when use_template is requested and the branded AiNxt template exists,
    # open it so slides inherit the master's branding (logo, colours, fonts) instead
    # of a blank white deck. Fall back to a fresh presentation if it can't be loaded.
    prs = None
    if use_template and os.path.exists(TEMPLATE_PATH):
        try:
            prs = _Presentation(TEMPLATE_PATH)
            # Drop any placeholder slides shipped inside the template so we start clean.
            _xml_slides = prs.slides._sldIdLst
            for _sld in list(_xml_slides):
                _xml_slides.remove(_sld)
            logger.info(f"generate_pptx: using branded template {TEMPLATE_PATH}")
        except Exception as exc:
            logger.warning(f"generate_pptx: template load failed, using blank deck: {exc}")
            prs = None
    if prs is None:
        prs = _Presentation()
    prs.slide_width  = _Inches(10)
    prs.slide_height = _Inches(5.625)
    blank_layout = prs.slide_layouts[6]  # blank — draw everything manually

    # Resolve theme colours; fall back to dark_executive defaults
    theme_cfg = PPTX_THEMES.get(theme) or PPTX_THEMES["dark_executive"]

    is_rich = bool(sections and isinstance(sections[0], dict) and "slide_type" in sections[0])
    rich = sections if is_rich else _legacy_to_rich(title, sections)

    for data in rich:
        slide  = prs.slides.add_slide(blank_layout)
        stype  = (data.get("slide_type") or "content").lower().strip()
        renderer = _SLIDE_RENDERERS.get(stype, _render_content)
        renderer(slide, data, title, prs, theme_cfg)

        notes = (data.get("speaker_notes") or "").strip()
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PDF ──────────────────────────────────────────────────────────────────────
#
# Premium branded executive PDF generator.
# Library: reportlab  (pip install reportlab)
#
# Design system mirrors generate_docx() exactly:
#   - Branded cover page with topic-specific colour palette (domain-driven)
#   - Running header: org name | doc title | page number  (accent underline)
#   - Footer: "Confidential — Internal Use Only"  (accent top rule)
#   - H1: full-width primary-colour band, white bold 16pt + thick accent stripe
#   - H2: primary-colour bold 13pt, accent underline rule
#   - H3: accent-colour italic 11pt
#   - Body: Helvetica 10.5pt, 1.15 line spacing, justified
#   - Callout boxes: badge (accent bg, white text) + insight (light_bg)
#   - Data tables: primary header row, alternating light_bg/white rows
#   - Bullet lists with → arrow markers in accent colour
#   - Section dividers between H1 sections

def _hex_to_rgb(hex6: str) -> tuple:
    """Convert 6-char hex string (no #) to (R, G, B) int tuple."""
    h = hex6.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _rl_color(hex6: str):
    """Return a ReportLab Color from a 6-char hex string (no #)."""
    try:
        from reportlab.lib.colors import Color
        r, g, b = _hex_to_rgb(hex6)
        return Color(r / 255, g / 255, b / 255)
    except Exception:
        from reportlab.lib.colors import black
        return black


def generate_pdf(title: str, sections: list, domain: str | None = None) -> bytes:
    """
    Generate a premium branded executive PDF report using ReportLab.

    Design system — Reference Standard:
      ┌─────────────────────────────────────────────────────────────────┐
      │  Cover page  — clean minimal:                                   │
      │    • Full-bleed primary top band (40 % of page)                 │
      │    • Large white title centred in band                          │
      │    • Accent stripe at band bottom                               │
      │    • Date + CONFIDENTIAL note below band on light-bg            │
      │    • Bottom primary bar                                         │
      ├─────────────────────────────────────────────────────────────────┤
      │  Body pages (no TOC — content starts immediately after cover)   │
      │    Header (single line, 3 zones):                               │
      │      LEFT:   [Doc Title] — [Short Title]                        │
      │      CENTER: Confidential | © 2025                              │
      │      RIGHT:  Page N                                             │
      │    Footer: accent rule + "Confidential — Internal Use Only"     │
      │    H1: full-width primary band + accent stripe, numbered (1.)   │
      │    H2: primary bold 13 pt + accent underline rule, numbered     │
      │    H3: accent italic 11 pt                                      │
      │    Body: Helvetica 10.5 pt, justified, 1.15 leading             │
      │    Callout: accent badge (descriptive label) + light-bg panel   │
      │    Table: primary header row, alternating light-bg/white rows   │
      │    Bullets: • bullet in accent colour                           │
      └─────────────────────────────────────────────────────────────────┘

    Section schema (same as generate_docx):
      {
        "heading":      str,
        "subheading":   str,          # H3 label
        "key_takeaway": str,          # → KEY callout
        "callout":      {"label": str, "text": str},
        "content":      str,          # body paragraphs (\\n\\n separated)
        "bullets":      list[str],
        "table":        {"headers": [...], "rows": [[...], ...]},
        "level":        int,          # 1 = H1 band, 2 = H2 (default)
        # legacy aliases still accepted:
        "verdict":      str,
        "recommendation": str,
        "key_insight":  str,
      }

    Library: reportlab  (pip install reportlab)
    """
    logger.info(
        f"[doc_generator] generate_pdf START | title={title!r} "
        f"sections={len(sections)} domain={domain!r}"
    )
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.colors import Color, white, black, HexColor
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY, TA_RIGHT
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            PageBreak, KeepTogether, HRFlowable,
        )
        from reportlab.platypus.flowables import Flowable
        from reportlab.pdfgen import canvas as _rl_canvas
    except ImportError:
        logger.error("[doc_generator] reportlab NOT installed")
        raise RuntimeError("reportlab not installed — run: pip install reportlab")

    import datetime as _dt

    # ── 1. Resolve domain palette ──────────────────────────────────────────────
    pal        = get_palette(domain)
    C_PRIMARY  = _rl_color(pal["primary"])
    C_ACCENT   = _rl_color(pal["accent"])
    C_LIGHT_BG = _rl_color(pal["light_bg"])
    C_BODY     = HexColor("#333333")
    C_WHITE    = white
    C_GREY     = HexColor("#888888")

    # ── 2. Page geometry ───────────────────────────────────────────────────────
    PAGE_W, PAGE_H = A4                    # 595.28 × 841.89 pt
    MARGIN         = 22 * mm              # ≈ 0.87 " — matches docx spec
    CONTENT_W      = PAGE_W - 2 * MARGIN
    HDR_H          = 14 * mm              # header band height
    FTR_H          = 14 * mm              # footer band height

    # ── 3. Header / Footer canvas callbacks ───────────────────────────────────
    #
    # Three-zone single-line header:
    #   LEFT:   [Doc Title] — [Short Title]
    #   CENTER: Confidential | © YYYY
    #   RIGHT:  Page N
    #
    _YEAR = _dt.date.today().year

    def _draw_header(c, page_num: int):
        """Draw the 3-zone running header on content pages (page_num >= 2)."""
        y_text = PAGE_H - 9 * mm

        # Thin primary top rule
        c.setStrokeColor(C_PRIMARY)
        c.setLineWidth(1.5)
        c.line(0, PAGE_H - HDR_H + 2, PAGE_W, PAGE_H - HDR_H + 2)

        # LEFT zone — doc title (truncated)
        short_title = title[:55] + ("\u2026" if len(title) > 55 else "")
        c.setFillColor(C_PRIMARY)
        c.setFont("Helvetica-Bold", 7.5)
        c.drawString(MARGIN, y_text, short_title)

        # CENTER zone — confidentiality + year
        conf_text = f"{_DOC_FOOTER_TEXT}  |  \u00a9 {_YEAR}"
        c.setFillColor(C_GREY)
        c.setFont("Helvetica", 7)
        conf_w = c.stringWidth(conf_text, "Helvetica", 7)
        c.drawString((PAGE_W - conf_w) / 2, y_text, conf_text)

        # RIGHT zone — page number
        pg_text = f"Page {page_num}"
        c.setFillColor(C_GREY)
        c.setFont("Helvetica", 7.5)
        pg_w = c.stringWidth(pg_text, "Helvetica", 7.5)
        c.drawString(PAGE_W - MARGIN - pg_w, y_text, pg_text)

        # Accent underline rule
        c.setStrokeColor(C_ACCENT)
        c.setLineWidth(0.75)
        c.line(MARGIN, PAGE_H - HDR_H, PAGE_W - MARGIN, PAGE_H - HDR_H)

    def _draw_footer(c):
        """Draw the footer rule + confidentiality notice."""
        c.setStrokeColor(C_ACCENT)
        c.setLineWidth(0.5)
        c.line(MARGIN, FTR_H, PAGE_W - MARGIN, FTR_H)
        c.setFillColor(C_GREY)
        c.setFont("Helvetica-Oblique", 7)
        conf_text = _DOC_FOOTER_TEXT
        conf_w = c.stringWidth(conf_text, "Helvetica-Oblique", 7)
        c.drawString((PAGE_W - conf_w) / 2, FTR_H - 5 * mm, conf_text)

    def _on_first_page(c, doc):
        """Page 1 = cover page — no header/footer."""
        c.saveState()
        c.restoreState()

    def _on_later_pages(c, doc):
        """Pages 2+ — draw header and footer."""
        c.saveState()
        _draw_header(c, doc.page)
        _draw_footer(c)
        c.restoreState()

    # ── 4. Custom Flowables ────────────────────────────────────────────────────

    class _H1Band(Flowable):
        """
        H1 heading — matches Word style exactly:
          - Bold primary-colour text (17 pt), left-aligned
          - Thin primary-colour bottom border rule beneath the text
          - No filled background band (clean text-only style)
        Supports word-wrapping for long headings.
        """

        _FONT_SIZE = 17
        _LINE_H    = 7 * mm    # line height per text row
        _RULE_GAP  = 2 * mm    # gap below last text line before rule
        _BOT_GAP   = 3 * mm    # gap below rule

        def __init__(self, text: str):
            super().__init__()
            self.text    = text
            self.width   = CONTENT_W
            self._lines  = None   # computed in wrap()
            self._height = None   # set in wrap()

        def _compute_lines(self, avail_w: float) -> list:
            """Word-wrap heading text to fit within available width."""
            # Approximate char width for Helvetica-Bold 17pt ≈ 10.2 pt/char
            max_w = avail_w
            words = self.text.split()
            lines, cur = [], ""
            for w in words:
                candidate = (cur + " " + w).strip()
                if len(candidate) * 10.2 <= max_w:
                    cur = candidate
                else:
                    if cur:
                        lines.append(cur)
                    cur = w
            if cur:
                lines.append(cur)
            return lines or [self.text]

        def wrap(self, availWidth, availHeight):
            self._lines  = self._compute_lines(availWidth)
            n            = len(self._lines)
            self._height = (n * self._LINE_H) + self._RULE_GAP + 1.5 + self._BOT_GAP
            return (availWidth, self._height)

        def draw(self):
            c     = self.canv
            lines = self._lines or [self.text]
            n     = len(lines)

            # rule sits at the bottom gap
            rule_y = self._BOT_GAP

            # Draw text lines from bottom up (ReportLab: y=0 is bottom)
            # First line baseline is just above the rule gap
            c.setFillColor(C_PRIMARY)
            c.setFont("Helvetica-Bold", self._FONT_SIZE)
            for i, line in enumerate(reversed(lines)):
                text_y = rule_y + self._RULE_GAP + (i * self._LINE_H)
                c.drawString(0, text_y, line)

            # Bottom border rule — thin primary-colour line (matching Word _add_h1)
            c.setStrokeColor(C_PRIMARY)
            c.setLineWidth(1.5)
            c.line(0, rule_y, self.width, rule_y)

    class _CalloutBox(Flowable):
        """
        Two-column callout box:
          [LABEL badge | accent fill, white text] [insight text | light_bg fill]
        The label is the callout's own descriptive label (e.g. 'Key Highlight',
        'Market Inflection', 'Strategic Vision') — not a fixed badge type.
        """

        def __init__(self, label: str, text: str):
            super().__init__()
            # Keep full descriptive label (no truncation) — badge auto-wraps to 2 lines
            # Guard against empty label — fall back to "Key Highlight"
            self.label   = (label.strip() or "Key Highlight").upper()
            self.text    = text
            self.width   = CONTENT_W
            self._height = None   # set in wrap()

        def _estimate_height(self, avail_w: float) -> float:
            badge_w        = 32 * mm   # wider badge to fit descriptive labels
            text_w         = avail_w - badge_w
            # Approximate character width at 9.5 pt Helvetica-Oblique ≈ 5.2 pt
            chars_per_line = max(1, int((text_w - 4 * mm) / 5.2))
            n_lines        = max(1, -(-len(self.text) // chars_per_line))  # ⌈ ⌉
            return max(12 * mm, n_lines * 5.2 * mm + 6 * mm)

        def wrap(self, availWidth, availHeight):
            self._height = self._estimate_height(availWidth)
            return (availWidth, self._height + 3 * mm)   # +3 mm bottom gap

        def draw(self):
            c       = self.canv
            h       = self._height or 12 * mm
            badge_w = 32 * mm   # wider to fit full descriptive labels
            text_w  = self.width - badge_w
            pad     = 2 * mm

            # Badge fill (accent colour — same as DOCX gold-border callout label)
            c.setFillColor(C_PRIMARY)
            c.rect(0, 0, badge_w, h, fill=1, stroke=0)
            # Insight panel fill — warm cream (#FFF8E7), matching DOCX callout box
            c.setFillColor(HexColor("#FFF8E7"))
            c.rect(badge_w, 0, text_w, h, fill=1, stroke=0)
            # Gold left accent border (matching DOCX gold border #E8B84B)
            c.setStrokeColor(HexColor("#E8B84B"))
            c.setLineWidth(2.5)
            c.line(0, 0, 0, h)
            # Thin grey border around the rest of the box
            c.setStrokeColor(HexColor("#DDDDDD"))
            c.setLineWidth(0.4)
            c.rect(0, 0, badge_w + text_w, h, fill=0, stroke=1)

            # Badge label — word-wrapped, vertically centred in badge
            # No character truncation — full descriptive label is shown
            c.setFillColor(C_WHITE)
            c.setFont("Helvetica-Bold", 7.5)
            words = self.label.split()
            badge_lines: list[str] = []
            cur_line = ""
            for w in words:
                candidate = (cur_line + " " + w).strip()
                if c.stringWidth(candidate, "Helvetica-Bold", 7.5) <= badge_w - 4 * mm:
                    cur_line = candidate
                else:
                    if cur_line:
                        badge_lines.append(cur_line)
                    cur_line = w
            if cur_line:
                badge_lines.append(cur_line)
            badge_lines = badge_lines[:3]  # max 3 lines for very long labels
            line_h_b = 9
            total_badge_h = len(badge_lines) * line_h_b
            start_y = h / 2 + total_badge_h / 2 - line_h_b
            for bl in badge_lines:
                lw = c.stringWidth(bl, "Helvetica-Bold", 7.5)
                c.drawString((badge_w - lw) / 2, start_y, bl)
                start_y -= line_h_b

            # Insight text — word-wrapped (body colour, italic, matching DOCX callout)
            c.setFillColor(C_BODY)
            c.setFont("Helvetica-Oblique", 9.5)
            max_line_w = text_w - 2 * pad
            words = self.text.split()
            lines: list[str] = []
            cur = ""
            for word in words:
                candidate = (cur + " " + word).strip()
                if c.stringWidth(candidate, "Helvetica-Oblique", 9.5) <= max_line_w:
                    cur = candidate
                else:
                    if cur:
                        lines.append(cur)
                    cur = word
            if cur:
                lines.append(cur)

            line_h = 4.2 * mm   # ~11.9 pt leading for 9.5 pt font
            # Start from top of box, leaving top padding
            y = h - pad - (9.5 / 2.835)   # top pad + approx ascender offset
            for line in lines:
                if y < pad:
                    break
                c.drawString(badge_w + pad, y, line)
                y -= line_h

    # ── 5. Paragraph styles ────────────────────────────────────────────────────
    _style_body = ParagraphStyle(
        "rl_body",
        fontName="Helvetica",
        fontSize=10.5,
        leading=13,
        textColor=C_BODY,
        alignment=TA_JUSTIFY,
        spaceAfter=8,
        spaceBefore=4,
    )
    _style_h2 = ParagraphStyle(
        "rl_h2",
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=16,
        textColor=C_PRIMARY,
        spaceBefore=14,
        spaceAfter=4,
    )
    _style_h3 = ParagraphStyle(
        "rl_h3",
        fontName="Helvetica-Oblique",
        fontSize=11,
        leading=14,
        textColor=C_ACCENT,
        spaceBefore=8,
        spaceAfter=2,
    )
    _style_bullet = ParagraphStyle(
        "rl_bullet",
        fontName="Helvetica",
        fontSize=10.5,
        leading=13,
        textColor=C_BODY,
        leftIndent=14,
        firstLineIndent=-14,
        spaceBefore=2,
        spaceAfter=2,
    )

    # ── 6. Helper: data table flowable ────────────────────────────────────────
    def _make_table(tbl_data: dict):
        headers   = tbl_data.get("headers") or []
        rows      = tbl_data.get("rows")    or []
        if not headers and not rows:
            return None
        col_count = max(len(headers), max((len(r) for r in rows), default=0))
        if col_count == 0:
            return None

        def _pad(row, n):
            return list(row[:n]) + [""] * max(0, n - len(row))

        table_data = []
        if headers:
            table_data.append([Paragraph(str(h), ParagraphStyle(
                "th", fontName="Helvetica-Bold", fontSize=9,
                textColor=C_WHITE, alignment=TA_CENTER,
            )) for h in _pad(headers, col_count)])
        for row in rows:
            table_data.append([Paragraph(str(cell), ParagraphStyle(
                "td", fontName="Helvetica", fontSize=9,
                textColor=C_BODY, alignment=TA_LEFT,
            )) for cell in _pad(row, col_count)])

        col_w = CONTENT_W / col_count
        tbl   = Table(table_data, colWidths=[col_w] * col_count,
                      repeatRows=1 if headers else 0)

        style_cmds = [
            ("GRID",          (0, 0), (-1, -1), 0.4, C_ACCENT),
            ("TOPPADDING",    (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("LEFTPADDING",   (0, 0), (-1, -1), 5),
            ("RIGHTPADDING",  (0, 0), (-1, -1), 5),
            ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ]
        if headers:
            style_cmds += [
                ("BACKGROUND", (0, 0), (-1, 0), C_PRIMARY),
                ("ALIGN",      (0, 0), (-1, 0), "CENTER"),
            ]
            for ri in range(1, len(table_data)):
                bg = C_LIGHT_BG if ri % 2 == 1 else C_WHITE
                style_cmds.append(("BACKGROUND", (0, ri), (-1, ri), bg))
        else:
            for ri in range(len(table_data)):
                bg = C_LIGHT_BG if ri % 2 == 0 else C_WHITE
                style_cmds.append(("BACKGROUND", (0, ri), (-1, ri), bg))

        tbl.setStyle(TableStyle(style_cmds))
        return tbl

    # ── 7. Helper flowables ────────────────────────────────────────────────────
    def _h2_rule():
        return HRFlowable(width="100%", thickness=1.5, color=C_ACCENT,
                          spaceAfter=6, spaceBefore=0)

    def _divider():
        return HRFlowable(width="100%", thickness=0.5, color=C_ACCENT,
                          spaceBefore=10, spaceAfter=10)

    # ── 8. Cover page flowable ─────────────────────────────────────────────────
    class _CoverPage(Flowable):
        """
        Clean minimal branded cover page.
        Layout:
          - Full-bleed primary top band (40 % of page height)
          - Large white title centred in band
          - Accent stripe at band bottom
          - Date below band on light-bg
          - CONFIDENTIAL note centred
          - Bottom primary bar
        """

        def wrap(self, availWidth, availHeight):
            return (availWidth, availHeight)

        def draw(self):
            c = self.canv

            # ── Geometry ─────────────────────────────────────────────────────
            # In flowable coords: (0,0) = bottom-left of content frame.
            # Content frame height = PAGE_H - topMargin - bottomMargin
            #   topMargin    = MARGIN + HDR_H
            #   bottomMargin = MARGIN + FTR_H
            frame_h = PAGE_H - (MARGIN + HDR_H) - (MARGIN + FTR_H)

            # ── Full-page white wash (bleeds into margins) ────────────────────
            c.setFillColor(C_WHITE)
            c.rect(
                -MARGIN, -(MARGIN + FTR_H),
                PAGE_W,  PAGE_H,
                fill=1, stroke=0,
            )

            # ── Primary top band (40 % of page height) ───────────────────────
            band_h   = PAGE_H * 0.40
            band_bot = frame_h - band_h   # y of band bottom in flowable coords

            c.setFillColor(C_PRIMARY)
            c.rect(-MARGIN, band_bot, PAGE_W, band_h, fill=1, stroke=0)

            # ── Accent stripe at band bottom ─────────────────────────────────
            c.setFillColor(C_ACCENT)
            c.rect(-MARGIN, band_bot, PAGE_W, 4, fill=1, stroke=0)

            # ── Document title — large, centred in band (word-wrapped) ───────
            c.setFillColor(C_WHITE)
            font_size   = 30
            max_title_w = PAGE_W - 2 * (MARGIN + 6 * mm)

            def _wrap_title(text, fs):
                words = text.split()
                lines, cur = [], ""
                for w in words:
                    candidate = (cur + " " + w).strip()
                    if c.stringWidth(candidate, "Helvetica-Bold", fs) <= max_title_w:
                        cur = candidate
                    else:
                        if cur:
                            lines.append(cur)
                        cur = w
                if cur:
                    lines.append(cur)
                return lines or [text]

            title_lines = _wrap_title(title, font_size)
            if len(title_lines) > 3:
                font_size   = 24
                title_lines = _wrap_title(title, font_size)
            if len(title_lines) > 4:
                font_size   = 20
                title_lines = _wrap_title(title, font_size)

            c.setFont("Helvetica-Bold", font_size)
            line_h_t  = font_size * 1.35
            total_th  = len(title_lines) * line_h_t
            # Vertically centre the title block in the band
            start_y_t = band_bot + (band_h + total_th) / 2 - line_h_t * 0.85
            for tl in title_lines:
                tw = c.stringWidth(tl, "Helvetica-Bold", font_size)
                c.drawString((CONTENT_W - tw) / 2, start_y_t, tl)
                start_y_t -= line_h_t

            # ── Light-bg panel below band ─────────────────────────────────────
            bottom_bar_h = 14 * mm
            c.setFillColor(C_LIGHT_BG)
            c.rect(
                -MARGIN, -(MARGIN + FTR_H) + bottom_bar_h,
                PAGE_W,  band_bot + (MARGIN + FTR_H) - bottom_bar_h,
                fill=1, stroke=0,
            )

            # ── Date line (centred, below band) ──────────────────────────────
            date_y    = band_bot - 32
            today_str = _dt.date.today().strftime("%B %Y")
            c.setFillColor(C_GREY)
            c.setFont("Helvetica-Oblique", 11)
            date_w = c.stringWidth(today_str, "Helvetica-Oblique", 11)
            c.drawString((CONTENT_W - date_w) / 2, date_y, today_str)

            # ── Thin accent divider ───────────────────────────────────────────
            c.setStrokeColor(C_ACCENT)
            c.setLineWidth(1)
            c.line(CONTENT_W * 0.25, date_y - 18,
                   CONTENT_W * 0.75, date_y - 18)

            # ── Confidential note ─────────────────────────────────────────────
            conf_text = _DOC_FOOTER_TEXT
            c.setFillColor(C_GREY)
            c.setFont("Helvetica-Oblique", 8)
            conf_w = c.stringWidth(conf_text, "Helvetica-Oblique", 8)
            c.drawString((CONTENT_W - conf_w) / 2, date_y - 38, conf_text)

            # ── Bottom primary bar ────────────────────────────────────────────
            c.setFillColor(C_PRIMARY)
            c.rect(
                -MARGIN, -(MARGIN + FTR_H),
                PAGE_W,  bottom_bar_h,
                fill=1, stroke=0,
            )
            # Accent stripe on top of bottom bar
            c.setFillColor(C_ACCENT)
            c.rect(
                -MARGIN, -(MARGIN + FTR_H) + bottom_bar_h - 3,
                PAGE_W,  3,
                fill=1, stroke=0,
            )

    # ── 9. Build story ─────────────────────────────────────────────────────────
    story: list = []

    # Cover page (full first page) — no TOC, content starts on page 2
    story.append(_CoverPage())
    story.append(PageBreak())

    # ── Section numbering counters ─────────────────────────────────────────────
    # H1 sections get sequential numbers: 1., 2., 3. …
    # H2 subsections get sub-numbers under the current H1: 1.1, 1.2 …
    _h1_counter = 0
    _h2_counter = 0

    # Body sections — direct, no TOC
    for sec_idx, sec in enumerate(sections):
        h              = (sec.get("heading")        or "").strip()
        subheading     = (sec.get("subheading")     or "").strip()
        key_takeaway   = (sec.get("key_takeaway")   or "").strip()
        callout        = sec.get("callout")
        verdict        = (sec.get("verdict")        or "").strip()
        recommendation = (sec.get("recommendation") or "").strip()
        key_insight    = (sec.get("key_insight")    or "").strip()
        content        = (sec.get("content")        or "").strip()
        bullets        = sec.get("bullets") or []
        table          = sec.get("table")
        level          = int(sec.get("level") or 2)

        block: list = []

        # Strip any leading number/bullet already in the heading text (e.g. "1)", "1.", "1 ")
        # to avoid double-numbering like "1.  1) Executive Summary".
        def _strip_num(text: str) -> str:
            return re.sub(r"^\s*\d+[\.\)]\s*", "", text).strip()

        # Heading with section numbering
        if h:
            if level == 1:
                _h1_counter += 1
                _h2_counter  = 0
                numbered_h   = f"{_h1_counter}.  {_strip_num(h)}"
                block.append(Spacer(1, 8 * mm))
                block.append(_H1Band(numbered_h))
                block.append(Spacer(1, 5 * mm))
            else:
                _h2_counter += 1
                numbered_h   = f"{_h1_counter}.{_h2_counter}  {_strip_num(h)}"
                block.append(Paragraph(numbered_h, _style_h2))
                block.append(_h2_rule())

        # Subheading (H3)
        if subheading:
            block.append(Paragraph(subheading, _style_h3))

        # Body paragraphs
        if content:
            for para_text in [p.strip() for p in content.split("\n\n") if p.strip()]:
                block.append(Paragraph(para_text, _style_body))

        # Callout / key-takeaway — use the callout's own descriptive label
        if callout and isinstance(callout, dict):
            block.append(Spacer(1, 3 * mm))
            block.append(_CalloutBox(
                callout.get("label") or "Key Highlight",
                callout.get("text") or "",
            ))
            block.append(Spacer(1, 3 * mm))
        elif key_takeaway:
            block.append(Spacer(1, 3 * mm))
            block.append(_CalloutBox("Key Highlight", key_takeaway))
            block.append(Spacer(1, 3 * mm))

        # Legacy inline callouts
        if key_insight:
            block.append(_CalloutBox("Key Insight", key_insight))
        if verdict:
            block.append(_CalloutBox("Verdict", verdict))
        if recommendation:
            block.append(_CalloutBox("Strategic Goal", recommendation))

        # Bullet list — • bullet in accent colour
        for b in bullets:
            if b and str(b).strip():
                bullet_line = (
                    f'<font color="#{pal["accent"]}"><b>\u2022</b></font>'
                    f'\u00a0\u00a0{str(b).strip()}'
                )
                block.append(Paragraph(bullet_line, _style_bullet))

        # Data table
        if table and isinstance(table, dict):
            tbl_flowable = _make_table(table)
            if tbl_flowable:
                block.append(Spacer(1, 4 * mm))
                block.append(tbl_flowable)
                block.append(Spacer(1, 4 * mm))

        # Section divider after H1 blocks (except the last section)
        if level == 1 and sec_idx < len(sections) - 1:
            block.append(_divider())

        story.extend(block)

    # ── 10. Render PDF ─────────────────────────────────────────────────────────
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf,
        pagesize=A4,
        leftMargin=MARGIN,
        rightMargin=MARGIN,
        topMargin=MARGIN + HDR_H,    # reserve space for header band
        bottomMargin=MARGIN + FTR_H, # reserve space for footer
        title=title,
        subject=title,
    )
    doc.build(story, onFirstPage=_on_first_page, onLaterPages=_on_later_pages)

    result_bytes = buf.getvalue()
    logger.info(
        f"[doc_generator] generate_pdf DONE | size={len(result_bytes):,} bytes "
        f"domain={domain!r} palette={pal['primary']} lib=reportlab"
    )
    return result_bytes



# ── XLSX ─────────────────────────────────────────────────────────────────────
#
# Matches the visual standard of the AiNxt_Payment_Report_2025.xlsx reference:
#   • FF-prefixed opaque colors (not 00-prefixed transparent)
#   • Channel-specific accent colors per sheet (green/orange/purple/red)
#   • KPI summary cards on the Executive Summary sheet
#   • Correct freeze panes: C5 on data sheets, C<first_data_row> on summary
#   • Channel row fills: E8F5E9 (UPI), FFF3E0 (Card), EDE7F6 (Intl)
#   • TwoCellAnchor charts on the summary sheet (LineChart + BarChart)
#   • Number formats: #,##0 integers, #,##0.0 one-decimal, ₹#,##0 currency
#   • Month/label column: center-aligned, bold
#   • Title cells: no border; data cells: thin CCCCCC all sides
# ─────────────────────────────────────────────────────────────────────────────

def generate_xlsx(title: str, sections: list) -> bytes:  # noqa: C901
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        raise RuntimeError("openpyxl not installed — run: pip install openpyxl")

    from datetime import datetime as _dt

    # ── Master colour palette ─────────────────────────────────────────────────
    # CRITICAL: ALL hex values must be 6-char with NO # prefix.
    # openpyxl stores colors as AARRGGBB internally; when you pass a 6-char hex
    # it prepends FF (fully opaque). Never pass 8-char values — that sets the
    # alpha channel and can produce transparent/invisible fills.
    PAL = {
        "header_dark":  "1A3C6E",   # title bars, main column headers
        "header_mid":   "2E75B6",   # section sub-headers
        "header_light": "D6E4F0",   # tertiary headers
        "upi":          "27AE60",   # UPI — green
        "card":         "E67E22",   # Card — amber
        "intl":         "8E44AD",   # International — purple
        "total":        "C0392B",   # Total / summary accent — red
        "upi_row":      "E8F5E9",   # pale green (UPI data rows)
        "card_row":     "FFF3E0",   # pale amber (Card data rows)
        "intl_row":     "EDE7F6",   # pale purple (Intl data rows)
        "total_row":    "FFFDE7",   # pale yellow (total rows)
        "row_white":    "FFFFFF",
        "row_alt":      "F5F9FF",   # near-white blue tint
        "text_dark":    "1A1A2E",
        "text_white":   "FFFFFF",
        "border_grey":  "CCCCCC",
        "border_dark":  "1A3C6E",
        "chart_upi":    "27AE60",
        "chart_card":   "E67E22",
        "chart_intl":   "8E44AD",
        "chart_line":   "2E75B6",
        "chart_total":  "C0392B",
    }

    # Channel accent colors — used for sheet-specific title/header fills
    # Maps a lowercase keyword found in the section heading to its accent color
    _CHANNEL_COLORS = {
        "upi":           PAL["upi"],
        "card":          PAL["card"],
        "international": PAL["intl"],
        "intl":          PAL["intl"],
        "remittance":    PAL["intl"],
        "comparison":    PAL["header_dark"],
        "channel":       PAL["header_dark"],
        "executive":     PAL["header_dark"],
        "summary":       PAL["header_dark"],
    }

    # Channel row-fill colors — pale tint matching the channel accent
    _CHANNEL_ROW_FILLS = {
        PAL["upi"]:          PAL["upi_row"],
        PAL["card"]:         PAL["card_row"],
        PAL["intl"]:         PAL["intl_row"],
        PAL["header_dark"]:  PAL["row_alt"],
        PAL["header_mid"]:   PAL["row_alt"],
    }

    _DARK_FILLS = {
        PAL["header_dark"], PAL["header_mid"],
        PAL["upi"], PAL["card"], PAL["intl"], PAL["total"],
    }

    def _accent_for(heading: str) -> str:
        """Return the channel accent color for a section heading."""
        h = heading.lower()
        for kw, color in _CHANNEL_COLORS.items():
            if kw in h:
                return color
        return PAL["header_dark"]

    def _row_fill_for(accent: str, row_index: int) -> str:
        """Return the alternating row fill for a given channel accent."""
        pale = _CHANNEL_ROW_FILLS.get(accent, PAL["row_alt"])
        return PAL["row_white"] if row_index % 2 == 0 else pale

    # ── Core style helpers ────────────────────────────────────────────────────

    def _fill(hex6: str) -> PatternFill:
        """Solid fill. hex6 must be exactly 6 chars, no # prefix."""
        return PatternFill(fill_type="solid", start_color=hex6, fgColor=hex6)

    def _font(color: str = "1A1A2E", size: int = 9,
              bold: bool = False, italic: bool = False) -> Font:
        return Font(name="Arial", color=color, size=size, bold=bold, italic=italic)

    def _align(horizontal: str = "center", wrap: bool = False) -> Alignment:
        return Alignment(horizontal=horizontal, vertical="center", wrap_text=wrap)

    def _border_std() -> Border:
        s = Side(style="thin", color=PAL["border_grey"])
        return Border(left=s, right=s, top=s, bottom=s)

    def _border_total() -> Border:
        thin  = Side(style="thin",   color=PAL["border_grey"])
        thick = Side(style="medium", color=PAL["border_dark"])
        return Border(left=thin, right=thin, top=thin, bottom=thick)

    def _no_border() -> Border:
        return Border()

    # ── Cell-level styling ────────────────────────────────────────────────────

    def style_title(cell, bg_hex: str, size: int = 14) -> None:
        """Title banner cell — no border (matches reference)."""
        cell.font      = _font(color=PAL["text_white"], size=size, bold=True)
        cell.fill      = _fill(bg_hex)
        cell.alignment = _align("center", wrap=True)
        cell.border    = _no_border()

    def style_col_header(cell, bg_hex: str, size: int = 9) -> None:
        """Column header cell — thin border, white text on dark fill."""
        text = PAL["text_white"] if bg_hex in _DARK_FILLS else PAL["text_dark"]
        cell.font      = _font(color=text, size=size, bold=True)
        cell.fill      = _fill(bg_hex)
        cell.alignment = _align("center", wrap=True)
        cell.border    = _border_std()

    def style_row_label(cell, bg_hex: str, size: int = 9) -> None:
        """Row label (first column) — left-aligned, bold, thin border."""
        cell.font      = _font(color=PAL["text_dark"], size=size, bold=True)
        cell.fill      = _fill(bg_hex)
        cell.alignment = _align("center", wrap=False)   # center matches reference
        cell.border    = _border_std()

    def style_data_num(cell, bg_hex: str, size: int = 9) -> None:
        """Numeric data cell — right-aligned, thin border."""
        cell.font      = _font(color=PAL["text_dark"], size=size, bold=False)
        cell.fill      = _fill(bg_hex)
        cell.alignment = _align("right", wrap=False)
        cell.border    = _border_std()

    def style_total_label(cell) -> None:
        cell.font      = _font(color=PAL["text_dark"], size=9, bold=True)
        cell.fill      = _fill(PAL["total_row"])
        cell.alignment = _align("center", wrap=False)
        cell.border    = _border_total()

    def style_total_num(cell) -> None:
        cell.font      = _font(color=PAL["text_dark"], size=9, bold=True)
        cell.fill      = _fill(PAL["total_row"])
        cell.alignment = _align("right", wrap=False)
        cell.border    = _border_total()

    def style_kpi_card(cell, bg_hex: str) -> None:
        """KPI summary card — large bold white text on channel color."""
        cell.font      = _font(color=PAL["text_white"], size=13, bold=True)
        cell.fill      = _fill(bg_hex)
        cell.alignment = _align("center", wrap=True)
        cell.border    = _no_border()

    def style_section_header(cell, bg_hex: str) -> None:
        """Section sub-header band (e.g. 'Monthly New User Enrollment…')."""
        cell.font      = _font(color=PAL["text_white"], size=10, bold=True)
        cell.fill      = _fill(bg_hex)
        cell.alignment = _align("left", wrap=False)
        cell.border    = _no_border()

    def style_subtitle(cell) -> None:
        cell.font      = _font(color="555555", size=8, italic=True)
        cell.fill      = _fill(PAL["row_white"])
        cell.alignment = _align("center", wrap=False)
        cell.border    = _no_border()

    # ── Numeric helpers ───────────────────────────────────────────────────────

    def _is_num(val) -> bool:
        cleaned = str(val).strip().replace(",", "").replace("₹", "").replace("Rs", "").rstrip("%")
        try:
            float(cleaned)
            return True
        except ValueError:
            return False

    def _to_num(val):
        """Convert string value to float; handle % by dividing by 100."""
        s = str(val).strip().replace(",", "").replace("₹", "").replace("Rs", "")
        is_pct = s.endswith("%")
        s = s.rstrip("%")
        try:
            n = float(s)
            return round(n / 100.0, 6) if is_pct else n
        except ValueError:
            return None

    def _nfmt(val) -> str:
        """Pick number format matching the reference: #,##0 / #,##0.0 / 0.0% / ₹#,##0."""
        s = str(val).strip()
        if s.endswith("%"):
            return "0.0%"
        if s.startswith("₹") or s.startswith("Rs"):
            return "₹#,##0"
        # One decimal if the value has a fractional part
        cleaned = s.replace(",", "")
        try:
            n = float(cleaned)
            if n != int(n):
                return "#,##0.0"
        except ValueError:
            pass
        return "#,##0"

    # ── Sheet-level helpers ───────────────────────────────────────────────────

    def _set_col_widths(ws, widths: dict) -> None:
        for col, w in widths.items():
            ws.column_dimensions[col].width = w

    def _set_row_height(ws, row: int, height: float) -> None:
        ws.row_dimensions[row].height = height

    # ── Build a full data sheet ───────────────────────────────────────────────

    def _build_data_sheet(ws, sheet_heading: str, doc_title: str,
                          headers: list, rows: list,
                          accent: str) -> int:
        """
        Build one data sheet matching the reference layout exactly.

        Layout:
          Row 1–2  : Merged title banner (accent color)
          Row 3    : Subtitle (reporting period / generated date)
          Row 4    : Column headers
          Row 5+   : Data rows (channel-tinted alternating fills)
          Last row : TOTAL row (FFFDE7 + medium bottom border)

        Returns the total_row number.
        """
        ws.sheet_view.showGridLines = False

        n_cols = len(headers)
        # col_letters: B = label col, C onwards = data cols
        col_letters = [get_column_letter(ci + 2) for ci in range(n_cols)]
        last_col    = col_letters[-1] if col_letters else "B"

        # Column widths — match reference: A=3, B=22-26, C=15-16, D+=13
        _set_col_widths(ws, {"A": 3, "B": 22})
        for col_l in col_letters[1:]:
            ws.column_dimensions[col_l].width = 13

        # ── Rows 1–2: Title banner ────────────────────────────────────────────
        ws.merge_cells(f"B1:{last_col}2")
        title_cell = ws["B1"]
        title_cell.value = sheet_heading
        style_title(title_cell, accent, size=14)
        _set_row_height(ws, 1, 27.75)
        _set_row_height(ws, 2, 15.0)

        # ── Row 3: Subtitle ───────────────────────────────────────────────────
        ws.merge_cells(f"B3:{last_col}3")
        sub = ws["B3"]
        sub.value = f"Reporting Period: {_dt.now().strftime('%B %Y')}  |  Source: AI Report Engine"
        style_subtitle(sub)
        _set_row_height(ws, 3, 13.5)

        # ── Row 4: Column headers ─────────────────────────────────────────────
        for ci, (col_l, hdr) in enumerate(zip(col_letters, headers)):
            c = ws[f"{col_l}4"]
            c.value = str(hdr)
            style_col_header(c, accent, size=9)
        _set_row_height(ws, 4, 21.75)

        # ── Detect numeric columns ────────────────────────────────────────────
        num_cols = [False] * n_cols
        for rv in rows:
            for ci, v in enumerate(rv):
                if ci < n_cols and _is_num(v):
                    num_cols[ci] = True

        # ── Data rows (start at row 5) ────────────────────────────────────────
        data_start = 5
        for ri, row_vals in enumerate(rows):
            er  = data_start + ri
            bg  = _row_fill_for(accent, ri)
            _set_row_height(ws, er, 19.5)

            for ci, (col_l, val) in enumerate(zip(col_letters, row_vals)):
                c      = ws[f"{col_l}{er}"]
                raw    = str(val) if val is not None else ""
                if num_cols[ci]:
                    parsed = _to_num(raw)
                    c.value        = parsed if parsed is not None else raw
                    c.number_format = _nfmt(raw)
                    style_data_num(c, bg)
                else:
                    c.value = raw
                    style_row_label(c, bg)

        data_end  = data_start + len(rows) - 1
        total_row = data_end + 1

        # ── Total row ─────────────────────────────────────────────────────────
        _set_row_height(ws, total_row, 19.5)
        for ci, col_l in enumerate(col_letters):
            c = ws[f"{col_l}{total_row}"]
            if ci == 0:
                c.value = "FY TOTAL"
                style_total_label(c)
            elif num_cols[ci]:
                c.value         = f"=SUM({col_l}{data_start}:{col_l}{data_end})"
                c.number_format = _nfmt(str(rows[0][ci]) if rows else "0")
                style_total_num(c)
            else:
                c.value = ""
                style_total_label(c)

        # ── Freeze panes: C5 (below header, right of label col) ──────────────
        ws.freeze_panes = "C5"

        return total_row

    # ── Build Executive Summary sheet ────────────────────────────────────────

    def _build_exec_summary(ws, doc_title: str, sections: list) -> None:
        """
        Executive Summary sheet matching the reference layout:
          Row 1–2  : Merged title banner (dark navy)
          Row 3    : Subtitle / reporting period
          Row 4    : Spacer
          Row 5–7  : KPI cards (4 colored merged cells side by side)
          Row 8    : Spacer
          Row 9    : (optional extra KPI row)
          Row 10   : Section sub-header band
          Row 11   : Column headers
          Row 12+  : Data rows
          Last     : Total row
        Then 2 charts below.
        """
        ws.sheet_view.showGridLines = False

        # Collect all table sections for the summary data table
        tbl_secs = [s for s in sections
                    if s.get("table") and isinstance(s["table"], dict)
                    and s["table"].get("headers") and s["table"].get("rows")]

        # Use the first table section's data for the main summary table
        main_sec     = tbl_secs[0] if tbl_secs else None
        main_headers = main_sec["table"]["headers"] if main_sec else []
        main_rows    = main_sec["table"]["rows"]    if main_sec else []

        n_data_cols  = len(main_headers) if main_headers else 4
        # Content columns: B through last data col
        col_letters  = [get_column_letter(ci + 2) for ci in range(n_data_cols)]
        last_col     = col_letters[-1] if col_letters else "I"

        # Column widths
        _set_col_widths(ws, {"A": 3, "B": 26})
        for col_l in col_letters[1:]:
            ws.column_dimensions[col_l].width = 13

        # ── Rows 1–2: Title banner ────────────────────────────────────────────
        ws.merge_cells(f"B1:{last_col}2")
        tc = ws["B1"]
        tc.value = doc_title
        style_title(tc, PAL["header_dark"], size=16)
        _set_row_height(ws, 1, 15.0)
        _set_row_height(ws, 2, 15.0)

        # ── Row 3: Subtitle ───────────────────────────────────────────────────
        ws.merge_cells(f"B3:{last_col}3")
        sub = ws["B3"]
        sub.value = (f"Reporting Period: January {_dt.now().year} – December {_dt.now().year}  |  "
                     f"Source: AI Report Engine  |  Segment: New Users")
        style_subtitle(sub)
        _set_row_height(ws, 3, 13.5)

        # ── Row 4: Spacer ─────────────────────────────────────────────────────
        _set_row_height(ws, 4, 6.0)

        # ── Rows 5–7: KPI cards ───────────────────────────────────────────────
        # Build 4 KPI cards from section summaries or bullet points
        kpi_data = _extract_kpis(sections)
        kpi_colors = [PAL["upi"], PAL["card"], PAL["intl"], PAL["total"]]

        # KPI cards span 2 columns each: B:C, D:E, F:G, H:I
        kpi_col_pairs = [("B", "C"), ("D", "E"), ("F", "G"), ("H", "I")]
        # Ensure we have enough columns
        all_kpi_cols = []
        for pair in kpi_col_pairs:
            all_kpi_cols.extend(pair)
        # Expand column widths for KPI area
        for col_l in all_kpi_cols:
            if col_l not in ("A", "B"):
                ws.column_dimensions[col_l].width = 13

        for ki, (kpi_label, kpi_value, kpi_sub) in enumerate(kpi_data[:4]):
            c1, c2 = kpi_col_pairs[ki]
            ws.merge_cells(f"{c1}5:{c2}7")
            kpi_cell = ws[f"{c1}5"]
            kpi_cell.value = f"{kpi_label}\n{kpi_value}"
            if kpi_sub:
                kpi_cell.value += f"\n{kpi_sub}"
            style_kpi_card(kpi_cell, kpi_colors[ki])

        _set_row_height(ws, 5, 18.0)
        _set_row_height(ws, 6, 31.5)
        _set_row_height(ws, 7, 21.75)

        # ── Row 8: Spacer ─────────────────────────────────────────────────────
        _set_row_height(ws, 8, 6.0)

        # ── Row 9: Spacer ─────────────────────────────────────────────────────
        _set_row_height(ws, 9, 6.0)

        # ── Row 10: Section sub-header band ───────────────────────────────────
        ws.merge_cells(f"B10:{last_col}10")
        sh = ws["B10"]
        sh.value = f"  {main_sec['heading'] if main_sec else 'Monthly Summary'} — Overview"
        style_section_header(sh, PAL["header_mid"])
        _set_row_height(ws, 10, 19.5)

        # ── Row 11: Column headers ────────────────────────────────────────────
        for ci, (col_l, hdr) in enumerate(zip(col_letters, main_headers)):
            c = ws[f"{col_l}11"]
            c.value = str(hdr)
            style_col_header(c, PAL["header_dark"], size=9)
        _set_row_height(ws, 11, 30.0)

        # ── Detect numeric columns ────────────────────────────────────────────
        num_cols = [False] * n_data_cols
        for rv in main_rows:
            for ci, v in enumerate(rv):
                if ci < n_data_cols and _is_num(v):
                    num_cols[ci] = True

        # ── Data rows (start at row 12) ───────────────────────────────────────
        data_start = 12
        for ri, row_vals in enumerate(main_rows):
            er = data_start + ri
            bg = PAL["row_white"] if ri % 2 == 0 else PAL["row_alt"]
            _set_row_height(ws, er, 15.0)

            for ci, (col_l, val) in enumerate(zip(col_letters, row_vals)):
                c   = ws[f"{col_l}{er}"]
                raw = str(val) if val is not None else ""
                if num_cols[ci]:
                    parsed = _to_num(raw)
                    c.value         = parsed if parsed is not None else raw
                    c.number_format = _nfmt(raw)
                    style_data_num(c, bg)
                else:
                    c.value = raw
                    style_row_label(c, bg)

        data_end  = data_start + len(main_rows) - 1
        total_row = data_end + 1

        # ── Total row ─────────────────────────────────────────────────────────
        _set_row_height(ws, total_row, 15.0)
        for ci, col_l in enumerate(col_letters):
            c = ws[f"{col_l}{total_row}"]
            if ci == 0:
                c.value = "FY TOTAL"
                style_total_label(c)
            elif num_cols[ci]:
                c.value         = f"=SUM({col_l}{data_start}:{col_l}{data_end})"
                c.number_format = _nfmt(str(main_rows[0][ci]) if main_rows else "0")
                style_total_num(c)
            else:
                c.value = ""
                style_total_label(c)

        # ── Freeze panes: C12 (matches reference) ────────────────────────────
        ws.freeze_panes = f"C{data_start}"

        # ── Charts below data ─────────────────────────────────────────────────
        if main_rows and len(main_rows) >= 4:
            _add_exec_charts(ws, col_letters, num_cols, main_headers,
                             data_start, data_end, total_row)

    def _extract_kpis(sections: list) -> list:
        """
        Extract 4 KPI tuples (label, value, sub) from sections.
        Uses bullet points or content snippets as KPI values.
        """
        kpis = []
        kpi_labels = [
            ("Total New Users\nEnrolled", "users"),
            ("Total Txn Volume\n(₹ Crore)", "volume"),
            ("Avg Txn per User", "avg"),
            ("Digital Adoption\nRate", "rate"),
        ]
        # Try to extract numeric highlights from bullets
        all_bullets = []
        for sec in sections:
            for b in (sec.get("bullets") or []):
                if b and str(b).strip():
                    all_bullets.append(str(b).strip())

        for i, (label, _) in enumerate(kpi_labels):
            if i < len(all_bullets):
                # Use bullet as the value line — truncate to 20 chars
                val = all_bullets[i][:30]
                sub = ""
            else:
                val = "—"
                sub = ""
            kpis.append((label, val, sub))
        return kpis

    def _add_exec_charts(ws, col_letters: list, num_cols: list,
                         headers: list, data_start: int, data_end: int,
                         total_row: int) -> None:
        """Add LineChart + BarChart to the Executive Summary sheet."""
        try:
            from openpyxl.chart import BarChart, LineChart, Reference
        except ImportError:
            return

        # Find numeric column indices (skip label col 0)
        num_indices = [ci for ci, is_n in enumerate(num_cols) if is_n and ci > 0]
        if not num_indices:
            return

        anchor_row = total_row + 2

        # ── Chart 1: LineChart for first numeric column ───────────────────────
        lc = LineChart()
        lc.type    = "line"
        lc.style   = 10
        lc.title   = f"{str(headers[num_indices[0]] if num_indices else 'Trend')} — {_dt.now().year}"
        lc.y_axis.title = str(headers[num_indices[0]]) if num_indices else "Value"
        lc.x_axis.title = "Month"
        lc.height  = 7.5
        lc.width   = 15

        lc_col = num_indices[0] + 2   # openpyxl 1-indexed, +2 for A-spacer + B-label
        lc_ref = Reference(ws, min_col=lc_col, max_col=lc_col,
                           min_row=11, max_row=data_end)
        lc.add_data(lc_ref, titles_from_data=True)
        lc_cats = Reference(ws, min_col=2, min_row=data_start, max_row=data_end)
        lc.set_categories(lc_cats)
        if lc.series:
            try:
                lc.series[0].graphicalProperties.line.solidFill = PAL["chart_line"]
                lc.series[0].graphicalProperties.line.width = 25000
                lc.series[0].smooth = True
            except Exception:
                pass
        ws.add_chart(lc, f"B{anchor_row}")

        # ── Chart 2: BarChart for remaining numeric columns ───────────────────
        if len(num_indices) >= 2:
            bc = BarChart()
            bc.type      = "col"
            bc.grouping  = "clustered"
            bc.style     = 10
            bc.title     = f"Channel-wise Transaction Volume — {_dt.now().year}"
            bc.y_axis.title = "Value"
            bc.x_axis.title = "Month"
            bc.height    = 7.5
            bc.width     = 15

            bc_min = num_indices[1] + 2
            bc_max = num_indices[-1] + 2
            bc_ref = Reference(ws, min_col=bc_min, max_col=bc_max,
                               min_row=11, max_row=data_end)
            bc.add_data(bc_ref, titles_from_data=True)
            bc_cats = Reference(ws, min_col=2, min_row=data_start, max_row=data_end)
            bc.set_categories(bc_cats)

            series_colors = [PAL["chart_upi"], PAL["chart_card"],
                             PAL["chart_intl"], PAL["chart_total"]]
            for si, series in enumerate(bc.series):
                try:
                    series.graphicalProperties.solidFill = series_colors[si % len(series_colors)]
                except Exception:
                    pass

            # Place bar chart to the right of the line chart
            bar_col = get_column_letter(2 + len(col_letters) // 2 + 1)
            ws.add_chart(bc, f"{bar_col}{anchor_row}")

    # ── Build Channel Comparison sheet ───────────────────────────────────────

    def _build_comparison_sheet(ws, doc_title: str, sections: list) -> None:
        """
        Channel Comparison sheet — metrics side by side per channel.
        Matches reference: B=metric label, C=UPI, D=Card, E=Intl, F=Total.
        """
        ws.sheet_view.showGridLines = False
        _set_col_widths(ws, {"A": 3, "B": 30, "C": 18, "D": 18, "E": 18, "F": 18})

        # Title
        ws.merge_cells("B1:F2")
        tc = ws["B1"]
        tc.value = f"Channel Comparison — {doc_title}"
        style_title(tc, PAL["header_dark"], size=13)
        _set_row_height(ws, 1, 31.5)
        _set_row_height(ws, 2, 15.0)

        # Subtitle
        ws.merge_cells("B3:F3")
        sub = ws["B3"]
        sub.value = f"Reporting Period: FY {_dt.now().year}  |  Source: AI Report Engine"
        style_subtitle(sub)
        _set_row_height(ws, 3, 13.5)

        # Column headers with channel colors
        ch_headers = ["Metric", "UPI", "Card Payment", "International", "Grand Total / Avg"]
        ch_colors  = [PAL["header_dark"], PAL["upi"], PAL["card"], PAL["intl"], PAL["total"]]
        for ci, (col_l, hdr, col_hex) in enumerate(
                zip(list("BCDEF"), ch_headers, ch_colors)):
            c = ws[f"{col_l}4"]
            c.value = hdr
            style_col_header(c, col_hex, size=10)
        _set_row_height(ws, 4, 27.75)

        # Gather comparison rows from sections
        comp_rows = _build_comparison_rows(sections)
        for ri, (label, upi_v, card_v, intl_v, total_v) in enumerate(comp_rows):
            er = 5 + ri
            bg = PAL["row_alt"] if ri % 2 == 0 else PAL["row_white"]
            _set_row_height(ws, er, 21.75)

            for ci, (col_l, val) in enumerate(
                    zip(list("BCDEF"), [label, upi_v, card_v, intl_v, total_v])):
                c = ws[f"{col_l}{er}"]
                c.value = val
                if ci == 0:
                    c.font      = _font(color=PAL["text_dark"], size=10, bold=True)
                    c.fill      = _fill(bg)
                    c.alignment = _align("left", wrap=False)
                    c.border    = _border_std()
                else:
                    c.font      = _font(color=PAL["text_dark"], size=10)
                    c.fill      = _fill(bg)
                    c.alignment = _align("center", wrap=False)
                    c.border    = _border_std()
                    if _is_num(str(val)):
                        parsed = _to_num(str(val))
                        if parsed is not None:
                            c.value         = parsed
                            c.number_format = _nfmt(str(val))

        ws.freeze_panes = None   # matches reference (no freeze on comparison sheet)

    def _build_comparison_rows(sections: list) -> list:
        """Build comparison rows from section data."""
        rows = []
        metrics = [
            "Total Volume (₹ Cr)",
            "Total Txn Count (Lakhs)",
            "New Users Activated",
            "Avg Txn Value (₹)",
            "YoY Growth (%)",
            "Digital Adoption Rate",
        ]
        # Try to extract values from table sections
        tbl_secs = [s for s in sections
                    if s.get("table") and s["table"].get("rows")]

        for metric in metrics:
            vals = ["—", "—", "—", "—"]
            rows.append([metric] + vals)
        return rows

    # ── Build workbook ────────────────────────────────────────────────────────

    wb = Workbook()

    # Separate sections
    tbl_sections = [s for s in sections
                    if s.get("table") and isinstance(s.get("table"), dict)
                    and s["table"].get("headers") and s["table"].get("rows")]

    # ── Sheet 1: Executive Summary ────────────────────────────────────────────
    ws_exec = wb.active
    ws_exec.title = "Executive Summary"
    _build_exec_summary(ws_exec, title, sections)

    # ── Sheets 2+: One per section with a table ───────────────────────────────
    used_names: set = {"Executive Summary", "Channel Comparison"}
    channel_sheets_built = 0

    for sec in tbl_sections:
        heading = (sec.get("heading") or "Data").strip()
        tbl     = sec["table"]
        headers = tbl.get("headers") or []
        rows    = tbl.get("rows") or []

        if not headers or not rows:
            continue

        # Derive sheet name — title-case, max 31 chars, unique
        # Strip characters forbidden by openpyxl: : [ ] * ? / \
        raw_name   = re.sub(r'[:\[\]*?/\\]', "", heading.title())[:28].strip() or "Sheet"
        sheet_name = raw_name
        suffix = 2
        while sheet_name in used_names:
            sheet_name = f"{raw_name[:25]} {suffix}"
            suffix += 1
        used_names.add(sheet_name)

        accent = _accent_for(heading)
        ws     = wb.create_sheet(title=sheet_name)
        _build_data_sheet(ws, heading, title, headers, rows, accent)
        channel_sheets_built += 1

    # ── Last sheet: Channel Comparison ───────────────────────────────────────
    ws_comp = wb.create_sheet(title="Channel Comparison")
    _build_comparison_sheet(ws_comp, title, sections)

    # ── Fallback: if no table sections, add a Key Points sheet ───────────────
    if not tbl_sections:
        all_bullets: list = []
        for sec in sections:
            h = (sec.get("heading") or "").strip()
            for b in (sec.get("bullets") or []):
                if b and str(b).strip():
                    all_bullets.append([h, str(b).strip()])
            h = ""
        if all_bullets:
            ws_kp = wb.create_sheet(title="Key Points")
            accent = PAL["header_dark"]
            _build_data_sheet(ws_kp, "Key Points", title,
                              ["Section", "Key Point"], all_bullets, accent)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── TXT ──────────────────────────────────────────────────────────────────────

def generate_txt(title: str, sections: list) -> bytes:
    lines = [title, "=" * len(title), ""]

    for sec in sections:
        h = (sec.get("heading") or "").strip()
        content = (sec.get("content") or "").strip()
        bullets = sec.get("bullets") or []

        if h:
            lines += [h, "-" * len(h)]
        if content:
            lines.append(content)
        for b in bullets:
            if b and str(b).strip():
                lines.append(f"  - {str(b).strip()}")
        lines.append("")

    return "\n".join(lines).encode("utf-8")


# ── CSV ──────────────────────────────────────────────────────────────────────

def generate_csv(title: str, sections: list) -> bytes:
    """Generate a UTF-8 CSV file from structured sections/tables."""
    import csv
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow([title])
    writer.writerow([])
    for sec in sections:
        heading = (sec.get("heading") or "").strip()
        tbl = sec.get("table")
        if tbl and isinstance(tbl, dict) and tbl.get("headers") and tbl.get("rows"):
            # Table section: heading row, then headers, then data rows
            if heading:
                writer.writerow([heading])
            writer.writerow(tbl["headers"])
            for row in tbl["rows"]:
                writer.writerow(row if isinstance(row, list) else [row])
            writer.writerow([])
        else:
            # Text section: heading + bullets + content
            if heading:
                writer.writerow([heading])
            for bullet in (sec.get("bullets") or []):
                if bullet and str(bullet).strip():
                    writer.writerow([str(bullet).strip()])
            content = (sec.get("content") or "").strip()
            if content:
                writer.writerow([content])
            writer.writerow([])
    return buf.getvalue().encode("utf-8")


# ── MARKDOWN ─────────────────────────────────────────────────────────────────

def generate_md(title: str, sections: list) -> bytes:
    lines = [f"# {title}", ""]

    for sec in sections:
        h = (sec.get("heading") or "").strip()
        content = (sec.get("content") or "").strip()
        bullets = sec.get("bullets") or []
        level = int(sec.get("level") or 2)

        hashes = "#" * min(level + 1, 6)
        if h:
            lines.append(f"{hashes} {h}")
        if content:
            lines.append("")
            lines.append(content)
        for b in bullets:
            if b and str(b).strip():
                lines.append(f"- {str(b).strip()}")
        lines.append("")

    return "\n".join(lines).encode("utf-8")


# ── DISPATCH ─────────────────────────────────────────────────────────────────

def generate(
    format_raw: str,
    title: str,
    sections: list,
    use_template: bool = False,
    theme: str = "dark_executive",
    domain: str | None = None,
) -> tuple:
    """
    Returns (bytes, ext, mime_type).
    format_raw: user-supplied format hint (e.g. "word", "pptx", "pdf").
    theme:  PPTX theme id — "dark_executive" | "light_modern" | "vibrant_tech"
    domain: industry/domain keyword used to select the colour palette for
            DOCX and PDF outputs (e.g. "payments", "ai", "healthcare").
            Resolved via get_palette(); falls back to "default" if unknown.
    """
    ext = FORMAT_EXTENSIONS.get((format_raw or "pdf").lower().strip(), "pdf")
    logger.info(f"[doc_generator] generate() | format={ext} title={title!r} domain={domain!r}")

    generators = {
        "docx": lambda: generate_docx(title, sections, domain=domain),
        "pptx": lambda: generate_pptx(title, sections, use_template=use_template, theme=theme),
        "pdf":  lambda: generate_pdf(title, sections, domain=domain),
        "xlsx": lambda: generate_xlsx(title, sections),
        "txt":  lambda: generate_txt(title, sections),
        "md":   lambda: generate_md(title, sections),
        "csv":  lambda: generate_csv(title, sections),
    }

    data = generators[ext]()
    mime = MIME_TYPES.get(ext, "application/octet-stream")
    return data, ext, mime


def text_to_sections(text: str) -> list:
    """
    Split parsed markdown / plain text into doc_generator sections.

    Splits on Markdown headings (# / ## / ###).  Each heading becomes a
    section dict with keys: heading, content, bullets, level.
    Falls back to a single section when no headings are found.
    """
    import re as _re
    sections: list = []
    parts = _re.split(r'\n(?=#{1,3} )', text.strip())
    for part in parts:
        part = part.strip()
        if not part:
            continue
        lines = part.split('\n')
        hm = _re.match(r'^(#{1,3})\s+(.+)', lines[0])
        if hm:
            level   = min(len(hm.group(1)), 2)   # 1 or 2
            heading = hm.group(2).strip()
            body    = '\n'.join(lines[1:]).strip()
        else:
            level   = 2
            heading = ""
            body    = part
        bullets: list       = []
        content_lines: list = []
        for line in body.split('\n'):
            s = line.strip()
            if s.startswith(('- ', '* ', '• ')):
                bullets.append(s.lstrip('-*• ').strip())
            else:
                content_lines.append(s)
        content = '\n\n'.join(
            p for p in '\n'.join(content_lines).split('\n\n') if p.strip()
        )
        sections.append({
            "heading": heading,
            "content": content,
            "bullets": bullets,
            "level":   level,
        })
    if not sections:
        sections = [{"heading": "", "content": text, "bullets": [], "level": 2}]
    return sections


def slugify(text: str) -> str:
    s = re.sub(r"[^\w\s-]", "", (text or "doc").lower())
    return re.sub(r"[\s]+", "_", s.strip())[:60] or "document"


# ── Explicit filename extraction ─────────────────────────────────────────────
# Patterns that match user phrases like:
#   "name should be env.pdf"   → "env"
#   "save as report.docx"      → "report"
#   "call it summary"          → "summary"
#   "filename: my-report"      → "my-report"
_EXPLICIT_FILENAME_PATTERNS = [
    # "name (it/the file/should be) <name>"
    r'(?:file\s*)?name\s+(?:it\s+|the\s+file\s+)?(?:should\s+be\s+|as\s+)?["\']?([A-Za-z0-9_\-\.]+)["\']?',
    # "save (it) as <name>"
    r'save\s+(?:it\s+)?as\s+["\']?([A-Za-z0-9_\-\.]+)["\']?',
    # "call it <name>"
    r'call\s+it\s+["\']?([A-Za-z0-9_\-\.]+)["\']?',
    # "filename: <name>" or "file name: <name>"
    r'file\s*name\s*[:\-]\s*["\']?([A-Za-z0-9_\-\.]+)["\']?',
    # "named <name>" or "titled <name>"
    r'(?:named|titled)\s+["\']?([A-Za-z0-9_\-\.]+)["\']?',
    # "output (file/name) <name>"
    r'output\s+(?:file\s+|name\s+)?["\']?([A-Za-z0-9_\-\.]+)["\']?',
]

# Format-related words to strip from an extracted explicit filename stem
_FORMAT_WORDS = frozenset({
    "pdf", "docx", "doc", "pptx", "ppt", "xlsx", "xls",
    "txt", "text", "md", "markdown", "word", "excel",
    "powerpoint", "presentation", "slides", "spreadsheet",
})

# Action words to strip when building a topic-based filename from the query
_ACTION_WORDS = frozenset({
    "generate", "create", "write", "make", "build", "produce",
    "give", "prepare", "draft", "compose", "design", "develop",
    "get", "show", "provide", "output", "export",
})


def extract_explicit_filename(query: str) -> str:
    """
    Scan *query* for an explicit filename instruction such as:
      "name should be env.pdf"  →  "env"
      "save as report.docx"     →  "report"
      "call it summary"         →  "summary"
      "filename: my-report"     →  "my-report"

    Returns the bare stem (no extension, no spaces) if found, else "".
    The caller is responsible for appending the correct extension.
    """
    lower = (query or "").lower()
    for pattern in _EXPLICIT_FILENAME_PATTERNS:
        m = re.search(pattern, lower)
        if m:
            raw = m.group(1).strip().strip("\"'")
            # Strip any file extension that was included in the match
            stem = re.sub(r"\.[a-z]{2,5}$", "", raw)
            # Strip format words that ended up as the whole stem
            if stem.lower() in _FORMAT_WORDS:
                continue
            # Sanitise: keep alphanumerics, hyphens, underscores
            stem = re.sub(r"[^\w\-]", "-", stem)
            stem = re.sub(r"[-_]+", "-", stem).strip("-_")
            if len(stem) >= 2:
                return stem.lower()
    return ""


def _topic_from_query(query: str, max_words: int = 6) -> str:
    """
    Extract a meaningful topic slug from a free-text query by:
      1. Removing action words (generate, create, write, …)
      2. Removing format words (pdf, docx, …)
      3. Taking up to *max_words* meaningful words
      4. Joining with hyphens, lowercased

    Returns "" if nothing meaningful remains.
    """
    words = re.findall(r"[A-Za-z0-9]+", query or "")
    stop = _ACTION_WORDS | _FORMAT_WORDS | {
        "me", "a", "an", "the", "in", "on", "for", "with",
        "using", "based", "from", "please", "can", "you",
        "about", "of", "and", "or", "to", "into", "as",
    }
    meaningful = [w for w in words if w.lower() not in stop][:max_words]
    if not meaningful:
        return ""
    return "-".join(w.lower() for w in meaningful)


# ── Output-type keyword map ───────────────────────────────────────────────────
# Maps intent keywords found in the user prompt → a short, human-readable label
# that is appended to the filename so the purpose is immediately obvious.
_OUTPUT_TYPE_KEYWORDS: list[tuple[list[str], str]] = [
    (["test case", "test cases", "testcase", "test suite", "test plan"],   "Test_Cases"),
    (["unit test", "unit tests"],                                           "Unit_Tests"),
    (["brd", "business requirement", "business requirements"],             "BRD"),
    (["frd", "functional requirement", "functional requirements"],         "FRD"),
    (["srs", "software requirement", "system requirement"],                "SRS"),
    (["user stor", "user stories"],                                        "User_Stories"),
    (["summary", "summarize", "summarise", "summarization"],               "Summary"),
    (["report", "analysis report", "analytical report"],                   "Report"),
    (["proposal", "project proposal"],                                     "Proposal"),
    (["roadmap", "product roadmap"],                                       "Roadmap"),
    (["presentation", "slide deck", "slides", "ppt"],                     "Presentation"),
    (["checklist"],                                                        "Checklist"),
    (["minutes", "meeting notes", "mom"],                                  "MoM"),
    (["sop", "standard operating procedure"],                              "SOP"),
    (["api doc", "api documentation", "swagger"],                          "API_Docs"),
    (["release note", "release notes", "changelog"],                       "Release_Notes"),
    (["audit", "audit report"],                                            "Audit_Report"),
    (["risk", "risk register", "risk assessment"],                         "Risk_Register"),
    (["gap analysis"],                                                     "Gap_Analysis"),
    (["comparison", "compare"],                                            "Comparison"),
]

# Common stop-words stripped from source document names before using them
# as filename prefixes (avoids "the_document_test_cases.xlsx" noise).
_DOC_STOP_WORDS = {
    "document", "doc", "file", "the", "a", "an", "of", "for", "and",
    "with", "using", "based", "on", "in", "to", "from", "by",
    "specification", "spec", "v1", "v2", "v3", "final", "draft",
    "copy", "new", "updated", "latest", "version",
}

_FILENAME_MAX_LEN = 80   # characters before the extension dot
_FILENAME_MIN_LEN = 3    # below this we fall back to "Document"


def _detect_output_type(prompt: str) -> str:
    """
    Scan the user prompt for known output-type keywords and return a
    short label (e.g. "Test_Cases", "Summary").  Returns "" if nothing matches.
    """
    lower = (prompt or "").lower()
    for keywords, label in _OUTPUT_TYPE_KEYWORDS:
        if any(kw in lower for kw in keywords):
            return label
    return ""


def _clean_source_name(raw_name: str) -> str:
    """
    Strip extension, stop-words, and special characters from an uploaded
    document name so it can be used as a clean filename prefix.

    Example:  "Payment_Specification_v2_FINAL.pdf"  →  "Payment_Specification"
    """
    import os as _os
    base = _os.path.splitext(raw_name or "")[0]          # drop extension
    base = re.sub(r"[^\w\s-]", " ", base)                # keep word chars
    tokens = re.split(r"[\s_\-]+", base)
    kept = [t for t in tokens if t.lower() not in _DOC_STOP_WORDS and len(t) > 1]
    result = "_".join(kept)
    return result[:40] or ""                              # cap prefix length


def smart_filename(
    *,
    title: str = "",
    question: str = "",
    source_doc_name: str = "",
    output_type_hint: str = "",
    fmt_ext: str = "",
) -> str:
    """
    Generate a context-aware, human-readable filename (without extension).

    Priority order:
      1. Explicit filename in query  — user said "name should be env.pdf" / "save as report"
         → extract stem, use as-is (lowercased, sanitised, hyphens)
      2. Source document + output type
         a. source_doc_name  (uploaded file, e.g. "Payment_Specification.pdf")
         b. output_type_hint OR auto-detected output type from question/title
      3. Topic extracted from query  — strip action/format words, take ≤6 meaningful words
      4. LLM-generated title         — slugified
      5. Fallback                    — "generated-document"

    Validation rules applied:
      - Only alphanumerics, underscores, hyphens (no spaces, no special chars)
      - Max {_FILENAME_MAX_LEN} characters (before extension)
      - Min {_FILENAME_MIN_LEN} characters (else fallback to "generated-document")
      - No leading/trailing underscores or hyphens
      - Consecutive separators collapsed to single hyphen

    Returns a slugified string ready to be combined with an extension:
        "rupay-specification-test-cases"
    """
    # ── Priority 1: Explicit filename in the user query ───────
    # Handles: "name should be env.pdf", "save as report.docx", "call it summary"
    _query = question or title
    if _query:
        explicit = extract_explicit_filename(_query)
        if explicit:
            # Validate & return immediately — user's intent is unambiguous
            explicit = re.sub(r"[^\w\-]", "-", explicit)
            explicit = re.sub(r"[-_]+", "-", explicit).strip("-_")
            if len(explicit) >= _FILENAME_MIN_LEN:
                return explicit[:_FILENAME_MAX_LEN]

    parts: list[str] = []

    # ── Priority 2a. Source document prefix ───────────────────
    if source_doc_name:
        clean = _clean_source_name(source_doc_name)
        if clean:
            parts.append(clean)

    # ── Priority 2b. Output type label ────────────────────────
    otype = output_type_hint or _detect_output_type(question) or _detect_output_type(title)
    if otype:
        parts.append(otype)

    # ── Priority 3. Topic from query (no source doc / output type yet) ──
    if not parts and question:
        topic = _topic_from_query(question, max_words=6)
        if topic:
            parts.append(topic)

    # ── Priority 4. LLM-generated title ───────────────────────
    if not parts and title:
        parts.append(slugify(title))

    # ── Priority 5. Hard fallback ──────────────────────────────
    if not parts:
        return "generated-document"

    # ── Assemble & sanitise ───────────────────────────────────
    raw = "-".join(p.strip("-_") for p in parts if p)
    raw = re.sub(r"[^\w\-]", "-", raw)          # replace non-word chars
    raw = re.sub(r"[-_]+", "-", raw)            # collapse consecutive separators
    raw = raw.strip("-_")                        # trim leading/trailing

    if len(raw) < _FILENAME_MIN_LEN:
        return "generated-document"

    return raw[:_FILENAME_MAX_LEN]


# ============================================================================
# MCP-server tool surface
# ----------------------------------------------------------------------------
# The functions below back the `doc_generator` MCP server
# (mcp/servers/doc_generator_server.py). They are intentionally small
# wrappers — every "render markdown to .docx / build a .pptx / write a .md"
# request from an agent flows through here so the outbox dir, filename
# sanitisation, and audit-friendly return shape (`{file: path}`) stay
# consistent across the whole MCP surface.
#
# These are deliberately separate from the (much larger) section-rendering
# pipeline higher up in this module: those tools are called from the
# Python-internal code path (workflow steps, n8n nodes); the MCP tools below
# are called from spec-compliant MCP clients (agent builder, ABStudio nodes,
# remote SSE consumers). Consolidating into a single file keeps the
# `doc_generator` concept undivided.
#
# Used by UC-71 (financial report), UC-82 (press release), UC-91 (deck
# generation), UC-93 (RFP response), UC-95 (policy / SOP drafting), UC-96
# (training material creation).
#
# Configuration (env vars):
#   DOC_GENERATOR_OUTPUT_DIR — where rendered files land
#                              (default ./outbox/mcp_outbox/generated_docs)
# ============================================================================

_DOC_GENERATOR_OUTPUT_DIR = os.getenv(
    "DOC_GENERATOR_OUTPUT_DIR",
    "./outbox/mcp_outbox/generated_docs",
)


def _doc_generator_out(name: str) -> str:
    """Return a safe absolute path under the doc_generator outbox for a
    user-supplied filename."""
    os.makedirs(_DOC_GENERATOR_OUTPUT_DIR, exist_ok=True)
    safe = re.sub(r"[^A-Za-z0-9._-]+", "_", name)
    return os.path.join(_DOC_GENERATOR_OUTPUT_DIR, safe)


def write_markdown(filename: str, content: str) -> dict:
    """Write markdown content to a .md file in the generated-docs outbox."""
    p = _doc_generator_out(filename if filename.endswith(".md") else filename + ".md")
    open(p, "w").write(content)
    return {"file": p}


def markdown_to_docx(filename: str, markdown_content: str, title: str = "") -> dict:
    """Render simple markdown (#/##/### headings, bullets, plain paragraphs)
    into a .docx file in the generated-docs outbox."""
    from docx import Document
    doc = Document()
    if title:
        doc.add_heading(title, level=0)
    for line in markdown_content.splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith("### "):
            doc.add_heading(s[4:], level=3)
        elif s.startswith("## "):
            doc.add_heading(s[3:], level=2)
        elif s.startswith("# "):
            doc.add_heading(s[2:], level=1)
        elif s.startswith(("- ", "* ")):
            doc.add_paragraph(s[2:], style="List Bullet")
        else:
            doc.add_paragraph(s)
    p = _doc_generator_out(filename if filename.endswith(".docx") else filename + ".docx")
    doc.save(p)
    return {"file": p}


def slides_to_pptx(filename: str, slides: list) -> dict:
    """Render slides into a .pptx in the generated-docs outbox.
    Each slide: {"title": str, "bullets": [str], "notes": str (optional)}."""
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation()
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s.get("title", "")
        body = slide.placeholders[1].text_frame
        bullets = s.get("bullets", [])
        if bullets:
            body.text = bullets[0]
            for b in bullets[1:]:
                para = body.add_paragraph()
                para.text = b
                para.font.size = Pt(18)
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
    p = _doc_generator_out(filename if filename.endswith(".pptx") else filename + ".pptx")
    prs.save(p)
    return {"file": p, "slides": len(slides)}
